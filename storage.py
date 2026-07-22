"""
客服考试结果看板 —— 可切换存储层
=====================================
本模块定义统一的存储接口 `BaseStorage`，并提供两种实现：

1. SQLiteStorage  —— 第一阶段（本地）实现，数据存于本地 SQLite 文件。
2. CloudStorage   —— 第二阶段（云端）占位实现，预留接口，演示如何从本地
                     无缝切换到云端（如自建 REST 服务 / Google Drive 等）。

切换方式：设置环境变量 `STORAGE_BACKEND=sqlite`（默认）或 `cloud`，
业务代码（server.py / 前端）无需任何改动即可迁移。

数据模型（Phase B 扩展）：
- cs_reps：增加 leave_date / status(active|left) / stage(新人|转正) / login_id
          以支持"客服从入职到离职的完整生命周期"。
- exam_sessions：增加 exam_type(onboarding|monthly|quarterly) / cycle_tag
          以支持周期考试（月/季考）与复用（同名考试多批次）。
- 新增 knowledge_dimensions（知识维度字典）、exam_question_dimensions
          （题→维度映射，按 exam_name+Q编号，同名批次自动复用）、accounts（登录）。
"""
import os
import json
import re
import sqlite3
import datetime
import threading
import hashlib
import hmac
from abc import ABC, abstractmethod
from typing import Optional

try:
    import psycopg2
    from psycopg2.extras import DictCursor
    _HAS_PG = True
except Exception:  # psycopg2 仅云端 Postgres 后端需要，本地 SQLite 可缺省
    _HAS_PG = False

# ----------------------------------------------------------------------------
# 配置（可通过环境变量覆盖，便于迁移到云端时只改配置、不改代码）
# ----------------------------------------------------------------------------
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DEFAULT_DB_PATH = os.path.join(BASE_DIR, "data", "exam.db")

DB_PATH = os.environ.get("STORAGE_SQLITE_PATH", DEFAULT_DB_PATH)
# 后端选择：sqlite（默认）| cloud
BACKEND = os.environ.get("STORAGE_BACKEND", "sqlite").lower()

# 默认知识维度字典（首次初始化时写入，可在管理端增删改；支持一题多维度自由组合）
DEFAULT_DIMS = [
    ("镜架", "Frames"),
    ("商城和促销", "Mall & Promotions"),
    ("镜片和处方", "Lenses & Prescription"),
    ("售中", "Mid-sales"),
    ("售后", "After-sales"),
]

SCHEMA = """
CREATE TABLE IF NOT EXISTS cs_reps (
    rep_id     TEXT PRIMARY KEY,
    name       TEXT NOT NULL,
    hire_date  TEXT,
    leave_date TEXT,
    status     TEXT DEFAULT 'active',   -- active / left
    stage      TEXT DEFAULT '新人',      -- 新人 / 转正
    login_id   TEXT,
    position   TEXT,                     -- 职级: Intern/Demoted/P1..P6/TL/QA
    channel    TEXT,                     -- 渠道: Email / LC / CC
    created_at TEXT DEFAULT CURRENT_TIMESTAMP
);

CREATE TABLE IF NOT EXISTS exam_sessions (
    session_id  INTEGER PRIMARY KEY AUTOINCREMENT,
    exam_name   TEXT NOT NULL,
    batch       TEXT NOT NULL,
    exam_date   TEXT NOT NULL,
    pass_score  REAL DEFAULT 60,
    exam_type   TEXT DEFAULT 'onboarding',  -- onboarding / monthly / quarterly
    cycle_tag   TEXT,
    note        TEXT,
    created_at  TEXT DEFAULT CURRENT_TIMESTAMP
);

CREATE TABLE IF NOT EXISTS exam_results (
    result_id  INTEGER PRIMARY KEY AUTOINCREMENT,
    session_id INTEGER NOT NULL,
    rep_id     TEXT NOT NULL,
    name       TEXT NOT NULL,
    subjects   TEXT,            -- JSON: {"Q1": 5, "Q2": 3, ...}
    total      REAL,
    passed     INTEGER,         -- 0 / 1
    created_at TEXT DEFAULT CURRENT_TIMESTAMP,
    FOREIGN KEY (session_id) REFERENCES exam_sessions(session_id),
    FOREIGN KEY (rep_id)     REFERENCES cs_reps(rep_id)
);

CREATE TABLE IF NOT EXISTS knowledge_dimensions (
    dim_id      INTEGER PRIMARY KEY AUTOINCREMENT,
    name_cn     TEXT NOT NULL,
    name_en     TEXT,
    description TEXT,
    created_at  TEXT DEFAULT CURRENT_TIMESTAMP
);

CREATE TABLE IF NOT EXISTS exam_question_dimensions (
    exam_name   TEXT NOT NULL,
    q_number    TEXT NOT NULL,
    dim_id      INTEGER NOT NULL,
    max_score   REAL,
    PRIMARY KEY (exam_name, q_number, dim_id),
    FOREIGN KEY (dim_id) REFERENCES knowledge_dimensions(dim_id)
);

CREATE TABLE IF NOT EXISTS exam_question_bank (
    exam_name   TEXT NOT NULL,
    q_number    TEXT NOT NULL,        -- 与 exam_results.subjects 的键（Q1..QN）对齐
    question_id INTEGER NOT NULL,
    seq         INTEGER,
    PRIMARY KEY (exam_name, q_number),
    FOREIGN KEY (question_id) REFERENCES questions(question_id)
);

-- 已上传考题 PPT 的元数据（供「考题 PPT 管理」列表/删除使用）
CREATE TABLE IF NOT EXISTS question_bank_meta (
    exam_name     TEXT PRIMARY KEY,   -- = 关联的考试名（与 exam_question_bank.exam_name 一致）
    orig_filename TEXT,               -- 上传时的真实 PPT 文件名
    uploaded_at   TEXT,              -- 最近一次导入时间
    q_count       INTEGER,            -- 导入题数
    dim_ids       TEXT                -- 导入时勾选的维度 ID（逗号分隔）
);

CREATE TABLE IF NOT EXISTS accounts (
    rep_id        TEXT PRIMARY KEY,
    login_id      TEXT UNIQUE,
    password_hash TEXT,
    role          TEXT DEFAULT 'csr',   -- admin / csr
    created_at    TEXT DEFAULT CURRENT_TIMESTAMP,
    FOREIGN KEY (rep_id) REFERENCES cs_reps(rep_id)
);

CREATE TABLE IF NOT EXISTS questions (
    question_id INTEGER PRIMARY KEY AUTOINCREMENT,
    q_type      TEXT NOT NULL DEFAULT 'single',  -- single / multiple / judge / essay
    category    TEXT,
    content     TEXT NOT NULL,
    options     TEXT,            -- JSON: [{"key":"A","text":"..."}, ...]
    answer      TEXT,            -- single:'A'; multiple:JSON['A','C']; judge:'true'/'false'; essay:NULL
    dim_id      INTEGER,
    score       REAL DEFAULT 5,
    explanation TEXT,
    source_exam TEXT,            -- 备注：这道题原始属于哪一场考试（用于按来源快速组卷）
    created_at  TEXT DEFAULT CURRENT_TIMESTAMP,
    FOREIGN KEY (dim_id) REFERENCES knowledge_dimensions(dim_id)
);

CREATE TABLE IF NOT EXISTS question_attachments (
    att_id      INTEGER PRIMARY KEY AUTOINCREMENT,
    question_id INTEGER NOT NULL,
    seq         INTEGER DEFAULT 0,
    filename    TEXT,            -- 原始文件名
    mime        TEXT,            -- image/png 等
    stored_path TEXT,            -- 已废弃（图片改存数据库后留空），保留以兼容旧行
    data        BLOB,            -- 图片二进制，直接存数据库，跨部署永久保存
    created_at  TEXT DEFAULT CURRENT_TIMESTAMP,
    FOREIGN KEY (question_id) REFERENCES questions(question_id)
);

CREATE TABLE IF NOT EXISTS exam_assignments (
    paper_id    INTEGER NOT NULL,
    rep_id      TEXT NOT NULL,
    created_at  TEXT DEFAULT CURRENT_TIMESTAMP,
    PRIMARY KEY (paper_id, rep_id),
    FOREIGN KEY (paper_id) REFERENCES exam_papers(paper_id),
    FOREIGN KEY (rep_id) REFERENCES cs_reps(rep_id)
);

CREATE TABLE IF NOT EXISTS exam_papers (
    paper_id    INTEGER PRIMARY KEY AUTOINCREMENT,
    title       TEXT NOT NULL,
    batch       TEXT DEFAULT 'online',
    exam_type   TEXT DEFAULT 'onboarding',
    status      TEXT DEFAULT 'draft',   -- draft / published / closed
    duration_min INTEGER DEFAULT 0,     -- 0 = 不限时
    open_at     TEXT,
    close_at    TEXT,
    pass_score  REAL DEFAULT 60,
    created_by  TEXT,
    created_at  TEXT DEFAULT CURRENT_TIMESTAMP
);

CREATE TABLE IF NOT EXISTS paper_questions (
    paper_id    INTEGER NOT NULL,
    question_id INTEGER NOT NULL,
    seq         INTEGER,
    score       REAL,
    PRIMARY KEY (paper_id, question_id),
    FOREIGN KEY (paper_id) REFERENCES exam_papers(paper_id),
    FOREIGN KEY (question_id) REFERENCES questions(question_id)
);

CREATE TABLE IF NOT EXISTS exam_attempts (
    attempt_id  INTEGER PRIMARY KEY AUTOINCREMENT,
    paper_id    INTEGER NOT NULL,
    rep_id      TEXT NOT NULL,
    start_time  TEXT,
    submit_time TEXT,
    status      TEXT DEFAULT 'in_progress',  -- in_progress / submitted / graded
    auto_score  REAL DEFAULT 0,
    manual_score REAL DEFAULT 0,
    total_score REAL DEFAULT 0,
    passed      INTEGER DEFAULT 0,
    created_at  TEXT DEFAULT CURRENT_TIMESTAMP,
    FOREIGN KEY (paper_id) REFERENCES exam_papers(paper_id),
    FOREIGN KEY (rep_id) REFERENCES cs_reps(rep_id)
);

CREATE TABLE IF NOT EXISTS exam_answers (
    attempt_id  INTEGER NOT NULL,
    question_id INTEGER NOT NULL,
    answer      TEXT,
    is_correct  INTEGER,     -- 客观题 0/1；essay 为 NULL（待人工判分）
    score       REAL DEFAULT 0,
    graded_by   TEXT,
    graded_at   TEXT,
    PRIMARY KEY (attempt_id, question_id),
    FOREIGN KEY (attempt_id) REFERENCES exam_attempts(attempt_id),
    FOREIGN KEY (question_id) REFERENCES questions(question_id)
);

CREATE TABLE IF NOT EXISTS system_config (
    key        TEXT PRIMARY KEY,
    value      TEXT,
    updated_at TEXT DEFAULT CURRENT_TIMESTAMP
);

CREATE INDEX IF NOT EXISTS idx_results_session ON exam_results(session_id);
CREATE INDEX IF NOT EXISTS idx_results_rep     ON exam_results(rep_id);
CREATE INDEX IF NOT EXISTS idx_paper_q_paper  ON paper_questions(paper_id);
CREATE INDEX IF NOT EXISTS idx_attempt_paper  ON exam_attempts(paper_id);
CREATE INDEX IF NOT EXISTS idx_attempt_rep    ON exam_attempts(rep_id);
"""


def _now() -> str:
    return datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")


# 密码哈希：固定 salt 的 sha256（内部培训工具，非高安全场景；上云后建议改为 OAuth + 服务端哈希）。
_PW_SALT = os.environ.get("PW_SALT", "exam-dashboard-pw-salt-v1")
def _hash_pw(plain: str) -> str:
    return hashlib.sha256((_PW_SALT + (plain or "")).encode("utf-8")).hexdigest()
def _verify_pw(plain: str, h: str) -> bool:
    return bool(h) and h == _hash_pw(plain)


def _compute_total_and_pass(subjects: dict, pass_score: float) -> tuple:
    """根据各项成绩计算总分与是否通过。"""
    if not subjects:
        return 0.0, 0
    total = sum(float(v) for v in subjects.values())
    passed = 1 if total >= float(pass_score) else 0
    return round(total, 2), passed


def normalize_options(raw):
    """把题库 options 统一成 [{key, text}] 标准结构。兼容三种历史/异常格式：
    - 标准对象数组: [{"key":"A","text":"..."}, ...]
    - 纯文本数组(legacy): ["A. 选项一", "B. 选项二", ...]
    - 字符串(可能是 JSON 或 "A. x|B. y" 这类拼接文本): 尝试解析
    返回 list[dict]；任何无法识别的内容都会被安全丢弃，绝不让调用方 500。
    这是数据层健壮性修复：早期测试/导入曾把 options 存成字符串数组，
    导致前端 optText 取不到选项文本（只显示"A"而非"A. optA"）。"""
    if not raw:
        return []
    arr = raw
    if isinstance(raw, str):
        try:
            arr = json.loads(raw)
        except (TypeError, ValueError):
            # 退化：以 "X. " 或 "X、" 分隔的纯文本
            arr = [p.strip() for p in re.split(r"(?<=[A-Za-z0-9])[.、]\s*", raw) if p.strip()]
    if not isinstance(arr, list):
        return []
    out = []
    for i, item in enumerate(arr):
        if isinstance(item, dict):
            key = str(item.get("key") or "").strip()
            text = str(item.get("text") or "").strip()
            if not key:
                key = chr(ord("A") + i)
            out.append({"key": key, "text": text})
        elif isinstance(item, str):
            m = re.match(r"^([A-Za-z0-9])[.、]\s*(.*)$", item.strip())
            if m:
                out.append({"key": m.group(1), "text": m.group(2).strip()})
            else:
                out.append({"key": chr(ord("A") + i), "text": item.strip()})
    return out


# ----------------------------------------------------------------------------
# 抽象接口
# ----------------------------------------------------------------------------
class BaseStorage(ABC):
    @abstractmethod
    def init(self) -> None: ...

    @abstractmethod
    def seed_if_empty(self) -> None: ...

    # ---- 人员 ----
    @abstractmethod
    def list_reps(self) -> list: ...

    @abstractmethod
    def get_rep_by_name(self, name: str): ...

    @abstractmethod
    def create_rep(self, data: dict) -> dict: ...

    @abstractmethod
    def update_rep(self, rep_id: str, data: dict) -> dict: ...

    @abstractmethod
    def delete_rep(self, rep_id: str) -> None: ...

    @abstractmethod
    def batch_update_reps(self, rep_ids: list, fields: dict) -> int: ...

    # ---- 考试批次 ----
    @abstractmethod
    def list_sessions(self) -> list: ...

    @abstractmethod
    def get_session(self, session_id: int) -> Optional[dict]: ...

    @abstractmethod
    def create_session(self, data: dict, results: list = None) -> dict: ...

    @abstractmethod
    def update_session(self, session_id: int, data: dict) -> dict: ...

    @abstractmethod
    def delete_session(self, session_id: int) -> None: ...

    # ---- 成绩 ----
    @abstractmethod
    def list_results(self, filters: dict = None) -> list: ...

    @abstractmethod
    def create_result(self, data: dict) -> dict: ...

    @abstractmethod
    def update_result(self, result_id: int, data: dict) -> dict: ...

    @abstractmethod
    def delete_result(self, result_id: int) -> None: ...

    @abstractmethod
    def reset_all_data(self, scope: str = "all") -> dict:
        """清空考试数据。scope:
        - "all": 客服/账号/批次/成绩/在线考试/题维度映射 全部清空（维度字典保留）
        - "results": 仅清空 批次+成绩+题维度映射（保留客服与维度字典）
        - "online": 仅清空在线考试（题/卷/作答/答案）三表，保留题维度映射
        """

    @abstractmethod
    def import_results_excel(self, path, exam_name=None, batch=None,
                             exam_date=None, pass_ratio=0.7,
                             full_score=None, pass_score=None,
                             orig_filename=None) -> dict: ...

    # ---- 知识维度 ----
    @abstractmethod
    def list_dimensions(self) -> list: ...

    @abstractmethod
    def create_dimension(self, data: dict) -> dict: ...

    @abstractmethod
    def update_dimension(self, dim_id: int, data: dict) -> dict: ...

    @abstractmethod
    def delete_dimension(self, dim_id: int) -> None: ...

    @abstractmethod
    def set_question_dimension(self, exam_name, q_number, dim_id, max_score=None) -> list: ...

    @abstractmethod
    def delete_question_dimension(self, exam_name, q_number, dim_id) -> None: ...

    @abstractmethod
    def get_exam_question_dimensions(self, exam_name) -> list: ...

    # ---- 维度分析 ----
    @abstractmethod
    def exam_dimension_distribution(self, exam_name) -> dict: ...

    @abstractmethod
    def rep_dimension_weakness(self, rep_id, session_ids=None) -> dict: ...

    # ---- 题库 / 在线考试 ----
    @abstractmethod
    def list_questions(self, filters: dict = None) -> list: ...

    @abstractmethod
    def create_question(self, data: dict) -> dict: ...

    @abstractmethod
    def update_question(self, qid: int, data: dict) -> dict: ...

    @abstractmethod
    def delete_question(self, qid: int) -> None: ...

    @abstractmethod
    def bulk_set_question_source(self, ids: list, source_exam: str) -> int: ...

    @abstractmethod
    def bulk_delete_questions(self, ids: list) -> int: ...

    @abstractmethod
    def import_questions_excel(self, path: str) -> dict: ...

    @abstractmethod
    def import_questions_ppt(self, path: str, exam_name: str,
                             dim_ids: list = None, orig_filename: str = None) -> dict: ...

    @abstractmethod
    def list_question_banks(self) -> list: ...

    @abstractmethod
    def delete_question_bank(self, exam_name: str) -> dict: ...

    @abstractmethod
    def get_question_by_exam_q(self, exam_name: str, q_number: str) -> Optional[dict]: ...

    @abstractmethod
    def export_questions_excel(self) -> str: ...

    @abstractmethod
    def list_papers(self, filters: dict = None) -> list: ...

    @abstractmethod
    def get_paper(self, paper_id: int) -> Optional[dict]: ...

    @abstractmethod
    def create_paper(self, data: dict) -> dict: ...

    @abstractmethod
    def update_paper(self, paper_id: int, data: dict) -> dict: ...

    @abstractmethod
    def delete_paper(self, paper_id: int) -> None: ...

    @abstractmethod
    def publish_paper(self, paper_id: int, open_at: str = None, close_at: str = None) -> dict: ...

    @abstractmethod
    def get_paper_questions(self, paper_id: int) -> list: ...

    @abstractmethod
    def set_paper_questions(self, paper_id: int, items: list) -> None: ...

    @abstractmethod
    def list_available_papers(self, rep_id: str) -> list: ...

    @abstractmethod
    def start_attempt(self, paper_id: int, rep_id: str) -> dict: ...

    @abstractmethod
    def get_attempt(self, attempt_id: int) -> dict: ...

    @abstractmethod
    def list_attempts(self, filters: dict = None) -> list: ...

    @abstractmethod
    def submit_attempt(self, attempt_id: int, answers: dict) -> dict: ...

    @abstractmethod
    def list_pending_grading(self) -> list: ...

    @abstractmethod
    def grade_essay(self, attempt_id: int, question_id: int, score: float, grader: str) -> dict: ...

    # ---- 看板聚合 ----
    @abstractmethod
    def individual_view(self, rep_id: str) -> dict: ...

    @abstractmethod
    def batch_view(self, session_id: int) -> dict: ...

    @abstractmethod
    def period_view(self, start: str, end: str) -> dict: ...

    @abstractmethod
    def overview(self) -> dict: ...


# ----------------------------------------------------------------------------
# 第一阶段：SQLite 实现
# ----------------------------------------------------------------------------
class SQLiteStorage(BaseStorage):
    def __init__(self, db_path: str = DB_PATH):
        self.db_path = db_path
        os.makedirs(os.path.dirname(db_path), exist_ok=True)
        # 连接策略：Flask 多线程下，共享单条连接并发访问会因 cursor 互相穿插而
        # 间歇性 500 / 返回串场错数据；多连接又易写锁争用。采用官方推荐做法——
        # 每请求一条独立连接（见 conn 属性，存入 flask.g），请求内不跨线程、
        # 请求间靠 SQLite 文件锁(busy_timeout)隔离。非请求上下文(init)退回线程局部。
        self._local = threading.local()

    def _new_conn(self):
        c = sqlite3.connect(self.db_path, check_same_thread=False)
        c.row_factory = sqlite3.Row
        c.execute("PRAGMA foreign_keys = ON")
        c.execute("PRAGMA busy_timeout = 5000")  # 并发写时等待而非直接报 locked
        return c

    @property
    def conn(self):
        # 优先使用当前请求内的连接（每请求一条，绝不跨线程共享）
        try:
            from flask import g
            if not hasattr(g, "db_conn") or g.db_conn is None:
                g.db_conn = self._new_conn()
            return g.db_conn
        except RuntimeError:
            # 不在 Flask 请求/应用上下文中（如启动时 init）
            t = self._local
            if not hasattr(t, "conn") or t.conn is None:
                t.conn = self._new_conn()
            return t.conn

    def init(self) -> None:
        self.conn.executescript(SCHEMA)
        self._migrate()                       # 兼容已存在的库（加列，数据零丢失）
        self._seed_default_dimensions()       # 默认维度字典（仅首次为空时写入）
        # 题目附件存储目录
        os.makedirs(os.path.join(BASE_DIR, "data", "attachments"), exist_ok=True)
        self.conn.commit()

    def _migrate(self) -> None:
        """对已有数据库做向前迁移：给旧表补上新列。"""
        cols = {r["name"] for r in self.conn.execute("PRAGMA table_info(cs_reps)")}
        for col, ddl in [
            ("leave_date", "ALTER TABLE cs_reps ADD COLUMN leave_date TEXT"),
            ("status", "ALTER TABLE cs_reps ADD COLUMN status TEXT DEFAULT 'active'"),
            ("stage", "ALTER TABLE cs_reps ADD COLUMN stage TEXT DEFAULT '新人'"),
            ("login_id", "ALTER TABLE cs_reps ADD COLUMN login_id TEXT"),
            ("position", "ALTER TABLE cs_reps ADD COLUMN position TEXT"),
            ("channel", "ALTER TABLE cs_reps ADD COLUMN channel TEXT"),
        ]:
            if col not in cols:
                self.conn.execute(ddl)
        cols2 = {r["name"] for r in self.conn.execute("PRAGMA table_info(exam_sessions)")}
        for col, ddl in [
            ("exam_type", "ALTER TABLE exam_sessions ADD COLUMN exam_type TEXT DEFAULT 'onboarding'"),
            ("cycle_tag", "ALTER TABLE exam_sessions ADD COLUMN cycle_tag TEXT"),
        ]:
            if col not in cols2:
                self.conn.execute(ddl)
        cols3 = {r["name"] for r in self.conn.execute("PRAGMA table_info(questions)")}
        if "source_exam" not in cols3:
            self.conn.execute("ALTER TABLE questions ADD COLUMN source_exam TEXT")

        # 图片改存数据库（跨部署永久保存）：旧库补 data 列，并清理只有磁盘路径的孤儿附件
        qa_cols = {r["name"] for r in self.conn.execute("PRAGMA table_info(question_attachments)")}
        if "data" not in qa_cols:
            self.conn.execute("ALTER TABLE question_attachments ADD COLUMN data BLOB")
        self.conn.execute("DELETE FROM question_attachments WHERE data IS NULL")

        # 增强3-需求1：得分率判定所需列
        for tbl in ("exam_results", "exam_attempts"):
            tcols = {r["name"] for r in self.conn.execute(f"PRAGMA table_info({tbl})")}
            if "score_rate" not in tcols:
                self.conn.execute(f"ALTER TABLE {tbl} ADD COLUMN score_rate REAL")
            if "full_score" not in tcols:
                self.conn.execute(f"ALTER TABLE {tbl} ADD COLUMN full_score REAL")
        # 系统设置键值表（及格线比例等，需求2 也复用）
        self.conn.execute(
            "CREATE TABLE IF NOT EXISTS system_config (key TEXT PRIMARY KEY, value TEXT, updated_at TEXT)")
        if not self._row(self.conn.execute("SELECT 1 FROM system_config WHERE key='pass_line_ratio'")):
            self.conn.execute(
                "INSERT INTO system_config (key, value, updated_at) VALUES ('pass_line_ratio','0.88',?)",
                (_now(),))
        # 增强3-需求2：积分系统新表 + 默认规则配置
        self.conn.execute("""CREATE TABLE IF NOT EXISTS points_account (
            rep_id TEXT PRIMARY KEY, total INTEGER DEFAULT 0, updated_at TEXT)""")
        self.conn.execute("""CREATE TABLE IF NOT EXISTS points_log (
            log_id INTEGER PRIMARY KEY AUTOINCREMENT, rep_id TEXT, rule_key TEXT,
            delta INTEGER, ref_type TEXT, ref_id TEXT, note TEXT, created_at TEXT,
            year INTEGER, quarter INTEGER)""")
        # 增强3-Q8：积分按年度/季度累计，跨年重新计。旧库补列 + 回填历史行的 year/quarter
        pl_cols = {r["name"] for r in self.conn.execute("PRAGMA table_info(points_log)")}
        if "year" not in pl_cols:
            self.conn.execute("ALTER TABLE points_log ADD COLUMN year INTEGER")
        if "quarter" not in pl_cols:
            self.conn.execute("ALTER TABLE points_log ADD COLUMN quarter INTEGER")
        self.conn.execute(
            "UPDATE points_log SET year=CAST(SUBSTR(created_at,1,4) AS INTEGER), "
            "quarter=CAST((CAST(SUBSTR(created_at,6,2) AS INTEGER)-1)/3 AS INTEGER)+1 "
            "WHERE year IS NULL AND created_at IS NOT NULL")
        self.conn.execute("""CREATE TABLE IF NOT EXISTS learning_materials (
            material_id INTEGER PRIMARY KEY AUTOINCREMENT, title TEXT, mtype TEXT,
            dim_id INTEGER, content TEXT, file_path TEXT, url TEXT, link_kind TEXT, created_at TEXT)""")
        # 旧库补列：link 类型资料的内容类型标签（word/excel/pdf/video/web）
        lm_cols = {r["name"] for r in self.conn.execute("PRAGMA table_info(learning_materials)")}
        if "link_kind" not in lm_cols:
            self.conn.execute("ALTER TABLE learning_materials ADD COLUMN link_kind TEXT")
        # 增强3-需求5：资料可标注多个知识维度（自由组合），与题目一题多维度保持一致
        self.conn.execute("""CREATE TABLE IF NOT EXISTS material_dimensions (
            material_id INTEGER NOT NULL, dim_id INTEGER NOT NULL,
            PRIMARY KEY (material_id, dim_id))""")
        self.conn.execute("""CREATE TABLE IF NOT EXISTS study_records (
            record_id INTEGER PRIMARY KEY AUTOINCREMENT, rep_id TEXT, material_id INTEGER,
            status TEXT DEFAULT 'opened', progress INTEGER DEFAULT 0, created_at TEXT,
            UNIQUE(rep_id, material_id))""")
        self.conn.execute("""CREATE TABLE IF NOT EXISTS mini_quiz (
            quiz_id INTEGER PRIMARY KEY AUTOINCREMENT, rep_id TEXT, dim_id INTEGER,
            question_ids TEXT, passed INTEGER DEFAULT 0, created_at TEXT)""")
        if not self._row(self.conn.execute("SELECT 1 FROM system_config WHERE key='points_rules'")):
            self.conn.execute(
                "INSERT INTO system_config (key, value, updated_at) VALUES ('points_rules',?,?)",
                (json.dumps({"participate": 10, "pass": {0.88: 20, 0.90: 30, 0.95: 40},
                             "material": 5, "mini_quiz": 10}, ensure_ascii=False), _now()))
        if not self._row(self.conn.execute("SELECT 1 FROM system_config WHERE key='points_threshold'")):
            self.conn.execute(
                "INSERT INTO system_config (key, value, updated_at) VALUES ('points_threshold','100',?)",
                (_now(),))
        # 增强6-需求Q6：周期目标（按月/按季）可配，默认按季、未启用(0)
        if not self._row(self.conn.execute("SELECT 1 FROM system_config WHERE key='points_period'")):
            self.conn.execute(
                "INSERT INTO system_config (key, value, updated_at) VALUES ('points_period','quarter',?)",
                (_now(),))
        if not self._row(self.conn.execute("SELECT 1 FROM system_config WHERE key='points_period_target'")):
            self.conn.execute(
                "INSERT INTO system_config (key, value, updated_at) VALUES ('points_period_target','0',?)",
                (_now(),))
        # 管理员密码（哈希存储）
        if not self._row(self.conn.execute("SELECT 1 FROM system_config WHERE key='admin_password_hash'")):
            import hashlib
            default_pw = "admin123"
            self.conn.execute(
                "INSERT INTO system_config (key, value, updated_at) VALUES ('admin_password_hash',?,?)",
                (hashlib.sha256(default_pw.encode()).hexdigest(), _now()))
        # 智能推荐：弱项维度数量(top_n)与小测题量(quiz_n)可配
        if not self._row(self.conn.execute("SELECT 1 FROM system_config WHERE key='recommend_top_n'")):
            self.conn.execute(
                "INSERT INTO system_config (key, value, updated_at) VALUES ('recommend_top_n','3',?)",
                (_now(),))
        if not self._row(self.conn.execute("SELECT 1 FROM system_config WHERE key='recommend_quiz_n'")):
            self.conn.execute(
                "INSERT INTO system_config (key, value, updated_at) VALUES ('recommend_quiz_n','5',?)",
                (_now(),))
        self.conn.commit()
        # 职级归一：Demoted / P1 合并为 "Demoted P1"（幂等，重复执行无害）
        self.conn.execute(
            "UPDATE cs_reps SET position='Demoted P1' WHERE position IN ('Demoted','P1')")
        self.conn.commit()
        # 增强3-需求5：维度拆分（一次性、幂等；上线前题目会清空，无需逐题重映射）
        self._migrate_dimension_split_v1()
        # 补考/单独开启：exam_assignments 增加每客服独立的开放/截止时间窗口
        self._migrate_assignment_window_v1()
        # 首次升级：回填历史得分率 + 按新标准重算 passed（失败不影响启动）
        try:
            self.recompute_score_rates()
        except Exception as _e:
            print(f"[migrate] recompute_score_rates skipped: {_e}")

    def _seed_default_dimensions(self) -> None:
        c = self._row(self.conn.execute("SELECT COUNT(*) c FROM knowledge_dimensions"))["c"]
        if c == 0:
            for cn, en in DEFAULT_DIMS:
                self.conn.execute(
                    "INSERT INTO knowledge_dimensions (name_cn, name_en, created_at) VALUES (?,?,?)",
                    (cn, en, _now()))
            self.conn.commit()

    def _migrate_dimension_split_v1(self) -> None:
        """增强3-需求5：将原有粗维度拆分为更细的维度（一次性、幂等）。
        - 商城和促销 → 商城 + 促销
        - 镜片和处方 → 镜片 + 处方
        - 售中 → 售中操作 + 物流 + 售中SOP
        - 售后 → 风险等级 + 售后SOP
        - 镜架 保留
        用「商城」是否已存在做幂等守卫；题目不需重映射（上线前清空重录），
        原维度改名后其下题目自然归入改名后的首个子维度。
        """
        if self._row(self.conn.execute("SELECT 1 FROM knowledge_dimensions WHERE name_cn='商城'")):
            return  # 已拆分过，跳过
        splits = {
            "商城和促销": [("商城", "Mall"), ("促销", "Promotions")],
            "镜片和处方": [("镜片", "Lenses"), ("处方", "Prescription")],
            "售中": [("售中操作", "Mid-sales Operations"), ("物流", "Logistics"), ("售中SOP", "Mid-sales SOP")],
            "售后": [("风险等级", "Risk Level"), ("售后SOP", "After-sales SOP")],
        }
        for old_cn, news in splits.items():
            old = self._row(self.conn.execute(
                "SELECT * FROM knowledge_dimensions WHERE name_cn=?", (old_cn,)))
            if not old:
                continue
            # 原维度改名为首个子维度（保留其 dim_id，题目映射不丢）
            self.conn.execute(
                "UPDATE knowledge_dimensions SET name_cn=?, name_en=? WHERE dim_id=?",
                (news[0][0], news[0][1], old["dim_id"]))
            # 其余子维度作为新行插入
            for cn, en in news[1:]:
                self.conn.execute(
                    "INSERT INTO knowledge_dimensions (name_cn, name_en, created_at) VALUES (?,?,?)",
                    (cn, en, _now()))
        self.conn.commit()

    def _migrate_assignment_window_v1(self) -> None:
        """补考/单独开启：为 exam_assignments 增加每客服独立的时间窗口（open_at/due_at）。"""
        cols = [r["name"] for r in self.conn.execute("PRAGMA table_info(exam_assignments)")]
        if "open_at" not in cols:
            self.conn.execute("ALTER TABLE exam_assignments ADD COLUMN open_at TEXT")
        if "due_at" not in cols:
            self.conn.execute("ALTER TABLE exam_assignments ADD COLUMN due_at TEXT")
        self.conn.commit()

    # ---- 补考 / 单独开启考试（每客服独立时间窗口） ----
    def set_assignment_window(self, paper_id, rep_id, open_at, due_at):
        """为某客服设置该试卷的独立开放/截止时间（补考）。空值表示沿用试卷全局窗口。"""
        rep_id = (rep_id or "").strip()
        if not rep_id:
            raise ValueError("rep_id 必填")
        self.conn.execute(
            "INSERT OR REPLACE INTO exam_assignments (paper_id, rep_id, open_at, due_at, created_at) "
            "VALUES (?,?,?,?,?)",
            (paper_id, rep_id, open_at or None, due_at or None, _now()))
        self.conn.commit()

    def remove_assignment_window(self, paper_id, rep_id):
        """撤销某客服的补考窗口，恢复为试卷全局规则。"""
        self.conn.execute(
            "UPDATE exam_assignments SET open_at=NULL, due_at=NULL WHERE paper_id=? AND rep_id=?",
            (paper_id, rep_id))
        self.conn.commit()

    def list_makeup_assignments(self, paper_id):
        """返回该试卷下设置了独立时间窗口的补考名单。"""
        return self._rows(self.conn.execute(
            "SELECT paper_id, rep_id, open_at, due_at, created_at FROM exam_assignments "
            "WHERE paper_id=? AND (open_at IS NOT NULL OR due_at IS NOT NULL) ORDER BY rep_id",
            (paper_id,)))
    def _row(self, cur) -> Optional[dict]:
        r = cur.fetchone()
        return dict(r) if r else None

    def _rows(self, cur) -> list:
        return [dict(r) for r in cur.fetchall()]

    # ---- 人员 ----
    def list_reps(self) -> list:
        cur = self.conn.execute(
            "SELECT * FROM cs_reps ORDER BY hire_date, rep_id"
        )
        return self._rows(cur)

    def get_rep_by_name(self, name: str):
        n = name.strip()
        rows = self._rows(self.conn.execute(
            "SELECT * FROM cs_reps WHERE LOWER(name) = LOWER(?) OR rep_id = ?",
            (n, n)))
        return rows[0] if rows else None

    def create_rep(self, data: dict) -> dict:
        rep_id = data["rep_id"].strip()
        self.conn.execute(
            "INSERT INTO cs_reps (rep_id, name, hire_date, leave_date, status, stage, login_id, position, channel, created_at) "
            "VALUES (?,?,?,?,?,?,?,?,?,?)",
            (rep_id, data["name"].strip(), data.get("hire_date"),
             data.get("leave_date"), data.get("status", "active"),
             data.get("stage", "新人"), data.get("login_id"),
             data.get("position"), data.get("channel"), _now()))
        self.conn.commit()
        # 同步建立登录账号（默认密码 = 工号，客服首次登录后可自行修改）
        self.ensure_account(rep_id)
        return self._row(self.conn.execute(
            "SELECT * FROM cs_reps WHERE rep_id=?", (rep_id,)))

    def update_rep(self, rep_id: str, data: dict) -> dict:
        fields, vals = [], []
        for k in ("name", "hire_date", "leave_date", "status", "stage", "login_id", "position", "channel"):
            if k in data:
                v = data[k]
                if isinstance(v, str):
                    v = v.strip()
                fields.append(f"{k}=?"); vals.append(v)
        if fields:
            vals.append(rep_id)
            self.conn.execute(
                f"UPDATE cs_reps SET {','.join(fields)} WHERE rep_id=?",
                vals)
            self.conn.commit()
        return self._row(self.conn.execute(
            "SELECT * FROM cs_reps WHERE rep_id=?", (rep_id,))) or {}

    def delete_rep(self, rep_id: str) -> None:
        # 同步清理该客服的成绩与考试会话（避免孤儿数据）
        self.conn.execute("DELETE FROM exam_answers WHERE attempt_id IN (SELECT attempt_id FROM exam_attempts WHERE rep_id=?)", (rep_id,))
        self.conn.execute("DELETE FROM exam_attempts WHERE rep_id=?", (rep_id,))
        self.conn.execute("DELETE FROM exam_results WHERE rep_id=?", (rep_id,))
        self.conn.execute("DELETE FROM accounts WHERE rep_id=?", (rep_id,))
        self.conn.execute("DELETE FROM cs_reps WHERE rep_id=?", (rep_id,))
        self.conn.commit()

    # ---- 账号 / 密码 ----
    def ensure_account(self, rep_id: str, default_pw: str = None):
        """确保存在登录账号；不存在则用默认密码（默认=工号）创建。返回账号行。"""
        row = self._row(self.conn.execute(
            "SELECT * FROM accounts WHERE rep_id=?", (rep_id,)))
        if row:
            return row
        pw = default_pw if default_pw is not None else rep_id
        self.conn.execute(
            "INSERT INTO accounts (rep_id, login_id, password_hash, role) VALUES (?,?,?,?)",
            (rep_id, rep_id, _hash_pw(pw), "csr"))
        self.conn.commit()
        return self._row(self.conn.execute(
            "SELECT * FROM accounts WHERE rep_id=?", (rep_id,)))

    def verify_rep_password(self, rep_id: str, plain: str) -> bool:
        self.ensure_account(rep_id)
        row = self._row(self.conn.execute(
            "SELECT password_hash FROM accounts WHERE rep_id=?", (rep_id,)))
        return bool(row) and _verify_pw(plain, row["password_hash"])

    def set_rep_password(self, rep_id: str, plain: str) -> None:
        self.ensure_account(rep_id)
        self.conn.execute(
            "UPDATE accounts SET password_hash=? WHERE rep_id=?",
            (_hash_pw(plain), rep_id))
        self.conn.commit()

    def reset_rep_password(self, rep_id: str, default_pw: str = None) -> str:
        """管理员重置：密码回到默认（默认=工号），返回明文默认密码。"""
        pw = default_pw if default_pw is not None else rep_id
        self.conn.execute(
            "INSERT OR REPLACE INTO accounts (rep_id, login_id, password_hash, role) "
            "VALUES (?,?,?,?)",
            (rep_id, rep_id, _hash_pw(pw), "csr"))
        self.conn.commit()
        return pw

    def batch_delete_reps(self, rep_ids: list) -> int:
        n = 0
        for rid in rep_ids:
            rid = (rid or "").strip()
            if not rid:
                continue
            # 先清所有引用 cs_reps(rep_id) 的子表，否则外键约束导致 500
            self.conn.execute("DELETE FROM exam_answers WHERE attempt_id IN (SELECT attempt_id FROM exam_attempts WHERE rep_id=?)", (rid,))
            self.conn.execute("DELETE FROM exam_attempts WHERE rep_id=?", (rid,))
            self.conn.execute("DELETE FROM exam_results WHERE rep_id=?", (rid,))
            self.conn.execute("DELETE FROM exam_assignments WHERE rep_id=?", (rid,))
            self.conn.execute("DELETE FROM accounts WHERE rep_id=?", (rid,))
            # 无外键回指的孤儿记录一并清理，保持数据干净
            self.conn.execute("DELETE FROM mini_quiz WHERE rep_id=?", (rid,))
            self.conn.execute("DELETE FROM study_records WHERE rep_id=?", (rid,))
            self.conn.execute("DELETE FROM points_log WHERE rep_id=?", (rid,))
            self.conn.execute("DELETE FROM points_account WHERE rep_id=?", (rid,))
            self.conn.execute("DELETE FROM cs_reps WHERE rep_id=?", (rid,))
            n += 1
        self.conn.commit()
        return n

    def batch_update_reps(self, rep_ids: list, fields: dict) -> int:
        """批量更新客服字段（如在职状态 status、职级 position、渠道 channel）。
        仅允许白名单字段，避免任意列被误改。返回成功更新条数。"""
        allowed = {"status", "position", "channel"}
        sets, vals = [], []
        for k, v in (fields or {}).items():
            if k in allowed:
                sets.append(f"{k}=?")
                vals.append(None if v in (None, "") else str(v).strip())
        if not sets:
            return 0
        n = 0
        for rid in (rep_ids or []):
            rid = (rid or "").strip()
            if not rid:
                continue
            exist = self._row(self.conn.execute(
                "SELECT rep_id FROM cs_reps WHERE rep_id=?", (rid,)))
            if not exist:
                continue
            self.conn.execute(
                f"UPDATE cs_reps SET {', '.join(sets)} WHERE rep_id=?",
                vals + [rid])
            n += 1
        self.conn.commit()
        return n

    def import_reps(self, rows: list) -> dict:
        """批量导入客服。rows: list[dict]{rep_id,name,hire_date?,position?,channel?,status?}
        返回 {imported, errors:[{row,reason}]}。"""
        imported, errors = 0, []
        for i, r in enumerate(rows, start=2):  # Excel 第 2 行起（第 1 行表头）
            rep_id = (r.get("rep_id") or "").strip()
            name = (r.get("name") or "").strip()
            if not rep_id or not name:
                errors.append({"row": i, "reason": "工号与姓名必填"})
                continue
            try:
                self.conn.execute(
                    "INSERT OR REPLACE INTO cs_reps "
                    "(rep_id, name, hire_date, leave_date, status, stage, login_id, position, channel, created_at) "
                    "VALUES (?,?,?,?,?,?,?,?,?,?)",
                    (rep_id, name, r.get("hire_date") or None,
                     None, r.get("status") or "active", "新人",
                     None, r.get("position") or None, r.get("channel") or None, _now()))
                self.ensure_account(rep_id)
                imported += 1
            except Exception as e:
                errors.append({"row": i, "reason": str(e)})
        self.conn.commit()
        return {"imported": imported, "errors": errors}

    # ---- 考试批次 ----
    def list_sessions(self) -> list:
        cur = self.conn.execute(
            "SELECT s.*, "
            "(SELECT COUNT(*) FROM exam_results r WHERE r.session_id=s.session_id) AS result_count "
            "FROM exam_sessions s ORDER BY s.exam_date DESC, s.session_id DESC"
        )
        rows = self._rows(cur)
        # 补充统计
        out = []
        for row in rows:
            row["stats"] = self._session_stats(row["session_id"])
            out.append(row)
        return out

    def _session_stats(self, session_id: int) -> dict:
        cur = self.conn.execute(
            "SELECT COUNT(*) c, AVG(total) a, MIN(total) mn, MAX(total) mx, "
            "SUM(passed) p FROM exam_results WHERE session_id=?", (session_id,))
        r = self._row(cur)
        c = r["c"] or 0
        return {
            "count": c,
            "avg": round(r["a"], 2) if r["a"] is not None else None,
            "min": r["mn"],
            "max": r["mx"],
            "pass_count": r["p"] or 0,
            "pass_rate": round(r["p"] / c * 100, 1) if c else None,
        }

    def get_session(self, session_id: int) -> Optional[dict]:
        return self._row(self.conn.execute(
            "SELECT * FROM exam_sessions WHERE session_id=?", (session_id,)))

    # ---- 增强3-需求1：得分率判定 ----
    def get_config(self, key, default=None):
        r = self._row(self.conn.execute("SELECT value FROM system_config WHERE key=?", (key,)))
        return r["value"] if r else default

    def set_config(self, key, value):
        self.conn.execute(
            "INSERT INTO system_config (key, value, updated_at) VALUES (?,?,?) "
            "ON CONFLICT(key) DO UPDATE SET value=excluded.value, updated_at=excluded.updated_at",
            (key, str(value), _now()))
        self.conn.commit()

    def get_admin_password_hash(self) -> str:
        """获取管理员密码的哈希值（SHA256）。"""
        return self.get_config("admin_password_hash", "")

    def update_admin_password(self, new_password: str):
        """更新管理员密码（固定 salt 的 SHA256 哈希存储，与 seed 一致）。"""
        self.set_config("admin_password_hash", _hash_pw(new_password))

    def verify_admin_password(self, password: str) -> bool:
        """验证管理员密码（优先比对 salt 哈希；兼容旧版无 salt 哈希）。"""
        stored = self.get_admin_password_hash()
        if not stored:
            return False
        import hashlib
        if hmac.compare_digest(hashlib.sha256(password.encode()).hexdigest(), stored):
            return True
        return hmac.compare_digest(_hash_pw(password), stored)

    def get_pass_line_ratio(self) -> float:
        try:
            return float(self.get_config("pass_line_ratio", "0.88") or "0.88")
        except (TypeError, ValueError):
            return 0.88

    def _exam_full_score(self, exam_name):
        """从题库映射求该考试的满分（各小题分值之和）；无映射返回 None。"""
        if not exam_name:
            return None
        r = self._row(self.conn.execute(
            "SELECT SUM(q.score) s FROM exam_question_bank b "
            "JOIN questions q ON q.question_id=b.question_id WHERE b.exam_name=?",
            (exam_name,)))
        s = r["s"] if r else None
        return float(s) if s else None

    def _paper_full_score(self, paper_id):
        r = self._row(self.conn.execute(
            "SELECT SUM(score) s FROM paper_questions WHERE paper_id=?", (paper_id,)))
        s = r["s"] if r else None
        return float(s) if s else None

    def _score_rate_pass(self, total, full_score, pass_score, ratio, online=False):
        """返回 (score_rate, passed)。
        score_rate 仅用于「展示」：有满分则 = 总得分/满分。
        判定逻辑（与展示解耦，避免满分影响是否通过）：
        - 在线考试(online=True)：及格线 pass_score 解释为「百分比阈值」(如 60 = 60%)，
            passed = 总得分/满分*100 >= pass_score；pass_score 为空时退回全局比例*100。
        - 非在线(Excel/手动导入)：走「绝对分」判定 total >= pass_score
            （pass_score 为空时退回全局比例*100）。这样导入成绩填的"绝对及格线"是按
            绝对分比较，与满分无关；满分仅用于展示得分率(%)。
        """
        total = float(total or 0)
        sr = round(total / full_score, 4) if (full_score and full_score > 0) else None
        if online:
            ps = float(pass_score) if pass_score is not None else (float(ratio) * 100 if ratio is not None else 60.0)
            return sr, (1 if (sr is not None and sr * 100 >= ps) else 0)
        ps = float(pass_score) if pass_score is not None else (float(ratio) * 100 if ratio is not None else 60.0)
        return sr, (1 if total >= ps else 0)

    def create_session(self, data: dict, results: list = None) -> dict:
        cur = self.conn.execute(
            "INSERT INTO exam_sessions (exam_name, batch, exam_date, pass_score, exam_type, cycle_tag, note, created_at) "
            "VALUES (?,?,?,?,?,?,?,?)",
            (data["exam_name"].strip(), data["batch"].strip(),
             data["exam_date"], float(data.get("pass_score", 60) or 60),
             data.get("exam_type", "onboarding"), data.get("cycle_tag"),
             data.get("note"), _now()))
        session_id = cur.lastrowid
        self.conn.commit()
        if results:
            for r in results:
                r["session_id"] = session_id
                self.create_result(r)
        return self.get_session(session_id)

    def update_session(self, session_id: int, data: dict) -> dict:
        fields, vals = [], []
        for k in ("exam_name", "batch", "exam_date", "note", "exam_type", "cycle_tag"):
            if k in data:
                v = data[k]
                if isinstance(v, str):
                    v = v.strip()
                fields.append(f"{k}=?"); vals.append(v)
        if "pass_score" in data:
            fields.append("pass_score=?"); vals.append(float(data["pass_score"] or 60))
        if fields:
            vals.append(session_id)
            self.conn.execute(
                f"UPDATE exam_sessions SET {','.join(fields)} WHERE session_id=?",
                vals)
            # 修改及格线/比例后，按得分率重算该批次所有成绩的是否通过
            pass_score = self.get_session(session_id)["pass_score"]
            ratio = self.get_pass_line_ratio()
            snote = self.get_session(session_id)["note"] or ""
            cur = self.conn.execute(
                "SELECT result_id, total, full_score FROM exam_results WHERE session_id=?",
                (session_id,))
            for r in cur.fetchall():
                if snote.startswith("online:"):
                    m = re.search(r"paper_id=(\d+)", snote)
                    pid = int(m.group(1)) if m else None
                    fs = r["full_score"] or (self._paper_full_score(pid) if pid else None)
                else:
                    # Excel/手动导入：保留导入时显式提供的满分（仅用于得分率展示），
                    # 不置空。绝对分判定与满分无关。
                    fs = r["full_score"]
                _, passed = self._score_rate_pass(r["total"], fs, pass_score, ratio,
                                                  online=snote.startswith("online:"))
                self.conn.execute(
                    "UPDATE exam_results SET passed=? WHERE result_id=?",
                    (passed, r["result_id"]))
            self.conn.commit()
        return self.get_session(session_id)

    def delete_session(self, session_id: int) -> None:
        self.conn.execute("DELETE FROM exam_results WHERE session_id=?", (session_id,))
        self.conn.execute("DELETE FROM exam_sessions WHERE session_id=?", (session_id,))
        self.conn.commit()

    def bulk_delete_sessions(self, ids: list) -> int:
        """批量删除考试批次（每批级联删其成绩）。返回成功删除的批次数。"""
        n = 0
        for sid in (ids or []):
            try:
                self.delete_session(int(sid))
                n += 1
            except Exception:
                pass
        return n

    # ---- 成绩 ----
    def list_results(self, filters: dict = None) -> list:
        filters = filters or {}
        sql = ("SELECT r.*, s.exam_name, s.batch, s.exam_date, s.pass_score, s.exam_type "
               "FROM exam_results r JOIN exam_sessions s ON r.session_id=s.session_id "
               "WHERE 1=1")
        vals = []
        if filters.get("rep_id"):
            sql += " AND r.rep_id=?"; vals.append(filters["rep_id"])
        if filters.get("session_id"):
            sql += " AND r.session_id=?"; vals.append(int(filters["session_id"]))
        if filters.get("start_date"):
            sql += " AND s.exam_date>=?"; vals.append(filters["start_date"])
        if filters.get("end_date"):
            sql += " AND s.exam_date<=?"; vals.append(filters["end_date"])
        if filters.get("exam_type"):
            sql += " AND s.exam_type=?"; vals.append(filters["exam_type"])
        if filters.get("passed") in ("0", "1", 0, 1):
            sql += " AND r.passed=?"; vals.append(int(filters["passed"]))
        sql += " ORDER BY s.exam_date, r.rep_id"
        cur = self.conn.execute(sql, vals)
        rows = self._rows(cur)
        for row in rows:
            row["subjects"] = json.loads(row["subjects"]) if row["subjects"] else {}
        return rows

    def create_result(self, data: dict) -> dict:
        session = self.get_session(int(data["session_id"]))
        pass_score = session["pass_score"] if session else 60
        ratio = self.get_pass_line_ratio()
        subjects = data.get("subjects") or {}
        if isinstance(subjects, str):
            subjects = json.loads(subjects)
        total = data.get("total")
        if total is None:
            total, _ = _compute_total_and_pass(subjects, pass_score)
        else:
            total = float(total)
        full_score = data.get("full_score")
        # 仅信任显式提供的满分。题库总和(_exam_full_score)只适用于在线考试
        # （total 即各小题分值之和）；Excel/手动导入的 total 与题库总和不在同一量纲，
        # 不可回退，否则得分率会被算成极小值（如 24/125≈0.19）。无满分则保持旧绝对分判定。
        score_rate, passed = self._score_rate_pass(
            total, full_score, pass_score, ratio,
            online=(session.get("note") or "").startswith("online:"))
        cur = self.conn.execute(
            "INSERT INTO exam_results (session_id, rep_id, name, subjects, total, passed, score_rate, full_score, created_at) "
            "VALUES (?,?,?,?,?,?,?,?,?)",
            (int(data["session_id"]), data["rep_id"].strip(), data["name"].strip(),
             json.dumps(subjects, ensure_ascii=False), total, passed, score_rate, full_score, _now()))
        self.conn.commit()
        # 增强3-需求2：通过考试仅给分层积分（不叠参与基础分），不通过无分
        try:
            rid = data["rep_id"].strip()
            if passed == 1:
                pd = self.pass_points_for(score_rate, passed=True)
                if pd:
                    self.award_points(rid, "pass", "result", cur.lastrowid, "通过考试", delta=pd)
        except Exception as _e:
            print(f"[points] create_result award skipped: {_e}")
        return self._row(self.conn.execute(
            "SELECT * FROM exam_results WHERE result_id=?", (cur.lastrowid,)))

    def update_result(self, result_id: int, data: dict) -> dict:
        existing = self._row(self.conn.execute(
            "SELECT * FROM exam_results WHERE result_id=?", (result_id,)))
        if not existing:
            return {}
        merged = dict(existing)
        merged.update(data)
        session = self.get_session(int(merged["session_id"]))
        pass_score = session["pass_score"] if session else 60
        ratio = self.get_pass_line_ratio()
        subjects = merged.get("subjects") or {}
        if isinstance(subjects, str):
            subjects = json.loads(subjects)
        total = merged.get("total")
        if total is None:
            total, _ = _compute_total_and_pass(subjects, pass_score)
        else:
            total = float(total)
        full_score = merged.get("full_score")
        if full_score is None:
            full_score = existing.get("full_score")
        # 不回退题库总和：仅信任已存/显式满分（同 create_result 说明）。
        score_rate, passed = self._score_rate_pass(
            total, full_score, pass_score, ratio,
            online=(session.get("note") or "").startswith("online:"))
        self.conn.execute(
            "UPDATE exam_results SET rep_id=?, name=?, subjects=?, total=?, passed=?, score_rate=?, full_score=? "
            "WHERE result_id=?",
            (merged["rep_id"].strip(), merged["name"].strip(),
             json.dumps(subjects, ensure_ascii=False), total, passed, score_rate, full_score, result_id))
        self.conn.commit()
        return self._row(self.conn.execute(
            "SELECT * FROM exam_results WHERE result_id=?", (result_id,)))

    def delete_result(self, result_id: int) -> None:
        self.conn.execute("DELETE FROM exam_results WHERE result_id=?", (result_id,))
        self.conn.commit()

    def bulk_delete_results(self, ids: list) -> int:
        """批量删除成绩记录。返回成功删除的条数。"""
        n = 0
        for rid in (ids or []):
            try:
                self.delete_result(int(rid))
                n += 1
            except Exception:
                pass
        return n

    def reset_all_data(self, scope: str = "all") -> dict:
        """清空考试数据。scope 见基类文档。
        返回 {scope, cleared:{表:删除行数}, preserved:[保留项说明]}。
        注意：exam_assignments 同时外键指向 exam_papers 与 cs_reps，必须在这两个父表
        之前删除，否则 cs_reps / exam_papers 会因 FK 约束删除失败被静默跳过。"""
        cur = self.conn.cursor()
        cleared = {}

        def _del(table, where=None, params=()):
            try:
                if where:
                    cur.execute(f"DELETE FROM {table} WHERE {where}", params)
                else:
                    cur.execute(f"DELETE FROM {table}")
                cleared[table] = cur.rowcount
            except Exception as e:  # 表可能不存在（旧库）或 FK 冲突
                cleared[table] = f"skip:{e}"

        if scope == "online":
            # 仅在线考试相关：先清子表（有外键指向父表），再 papers/questions
            _del("exam_assignments")
            _del("exam_answers")
            _del("exam_attempts")
            _del("paper_questions")
            _del("exam_question_bank")
            _del("question_bank_meta")
            _del("question_attachments")
            _del("exam_papers")
            _del("questions")
            preserved = ["knowledge_dimensions（维度字典）", "system_config（系统设置）",
                         "cs_reps（客服名单）", "accounts（登录账号）", "exam_results（成绩）",
                         "exam_sessions（考试批次）", "exam_question_dimensions（题-维度映射）",
                         "learning_materials（资料库）", "material_dimensions（资料维度）",
                         "points_account / points_log（积分）"]
        elif scope == "results":
            # 仅成绩相关，保留客服与维度字典及在线考试
            _del("exam_question_dimensions")
            _del("exam_results")
            _del("exam_sessions")
            preserved = ["knowledge_dimensions（维度字典）", "system_config（系统设置）",
                         "cs_reps（客服名单）", "accounts（登录账号）", "questions（题库）",
                         "exam_papers / exam_attempts（在线考试）", "exam_question_bank（考题PPT）",
                         "learning_materials（资料库）", "points_account / points_log（积分）",
                         "exam_assignments（考试分配）"]
        else:  # all
            # 先清子表（有外键指向父表），避免 FK 报错；
            # exam_assignments 须先于 papers/cs_reps；question_attachments / exam_question_bank 须先于 questions
            _del("exam_assignments")
            _del("exam_answers")
            _del("exam_attempts")
            _del("paper_questions")
            _del("exam_question_bank")
            _del("question_bank_meta")
            _del("question_attachments")
            _del("exam_papers")
            _del("questions")
            _del("exam_question_dimensions")
            _del("exam_results")
            _del("exam_sessions")
            _del("accounts")
            _del("cs_reps")
            preserved = ["knowledge_dimensions（维度字典）", "system_config（系统设置）",
                         "learning_materials（资料库）", "material_dimensions（资料维度）",
                         "points_account / points_log（积分）"]
        self.conn.commit()
        return {"scope": scope, "cleared": cleared, "preserved": preserved}

    def import_results_excel(self, path, exam_name=None, batch=None,
                             exam_date=None, pass_ratio=0.7,
                             full_score=None, pass_score=None,
                             orig_filename=None) -> dict:
        """解析问卷星式成绩 Excel（1 文件 = 1 批次），建 session + results。
        列: 序号|用户ID|提交答卷时间|所用时间|来源|来源详情|来自IP|总分|Name|Email|小题1..N
        - exam_name/batch 未传时从原始文件名 `名称_批次.xlsx` 推断
          （orig_filename 为网页上传时的真实文件名；path 可能是临时文件名，不能用来推断）
        - 及格线：pass_score > full_score×pass_ratio > 各场最高分×pass_ratio
        - 同人同场保留最高分；rep_id = 姓名归一化(upper+trim)
        返回 {session_id, exam_name, batch, exam_date, count, max_total, pass_score, q_count}
        """
        import os, re, datetime, openpyxl

        def _norm_id(name):
            return re.sub(r"\s+", " ", (name or "").strip()).upper()

        # 导入时跳过的非考生行：渠道标签 / 汇总行（问卷星导出常见 TOTAL/EMAIL/LC/CC/CSM 等）
        _JUNK_IDS = {"TOTAL", "EMAIL", "LC", "CC", "CSM", "CHANNEL", "CHANNELS",
                     "NAME", "AVG", "SUM", "GRAND", "OVERALL", "ALL",
                     "SCORE", "FINAL", "PASS", "FAILED"}

        def _parse_dt(s):
            if not s:
                return None
            s = str(s).strip()
            for fmt in ("%Y/%m/%d %H:%M:%S", "%Y-%m-%d %H:%M:%S", "%Y/%m/%d %H:%M",
                        "%Y-%m-%d %H:%M", "%Y/%m/%d", "%Y-%m-%d"):
                try:
                    return datetime.datetime.strptime(s, fmt)
                except ValueError:
                    continue
            return None

        def _isnum(v):
            return isinstance(v, (int, float)) and not isinstance(v, bool)

        # 用原始上传文件名推断考试名/批次（临时落盘路径不可靠）
        base = os.path.splitext(os.path.basename(orig_filename or path))[0]
        if not exam_name:
            if "_" in base:
                exam_name, b2 = base.rsplit("_", 1)
                batch = batch or b2
            else:
                exam_name = base
        exam_name = re.sub(r"\s+", " ", exam_name).strip()
        batch = (batch or "batch").strip()

        wb = openpyxl.load_workbook(path, data_only=True)
        ws = wb.active
        max_col = ws.max_column
        header = [ws.cell(1, c).value for c in range(1, max_col + 1)]
        idx_total = next((i for i, h in enumerate(header)
                          if h is not None and ("总分" in str(h)
                                                or "score" in str(h).lower())), None)
        idx_name = next((i for i, h in enumerate(header) if h and "Name" in str(h)), None)
        idx_date = next((i for i, h in enumerate(header)
                         if h is not None and ("提交答卷时间" in str(h)
                                               or "date" in str(h).lower())), None)
        if idx_total is None or idx_name is None:
            raise ValueError("Excel 缺少「总分」或「Name」列，无法识别（需问卷星式导出格式）")

        # 小题列识别：优先用「表头」判定（问卷星列头为 1..N / 小题N / QN），
        # 不再仅依赖「数字占比≥60%」——否则自由作答题等纯文本列会被整列丢弃，
        # 导致小题数少 1 且后续题号整体错位（PPT Q29 在成绩里找不到对应）。
        start_q = max(idx_name, idx_date if idx_date is not None else -1) + 2

        def _is_qcol_header(h):
            if h is None:
                return False
            s = str(h).strip()
            if re.match(r"^\d+$", s):                       # 纯数字表头 1..N
                return True
            if re.search(r"小题|题目|question", s, re.I):   # 小题N / 题目 / question
                return True
            if re.match(r"^Q\d", s, re.I):                  # Q1..QN
                return True
            return False

        q_cols = [c for c in range(start_q, max_col + 1) if _is_qcol_header(header[c - 1])]
        if not q_cols:
            # 兜底：非标准表头时退回数字占比法
            for c in range(start_q, max_col + 1):
                nums = tot = 0
                for r in range(2, min(ws.max_row, 60) + 1):
                    v = ws.cell(r, c).value
                    if v is None or str(v).strip() == "":
                        continue
                    tot += 1
                    if _isnum(v):
                        nums += 1
                if tot and nums / tot >= 0.6:
                    q_cols.append(c)

        rows, dates, max_total = [], [], 0
        for r in range(2, ws.max_row + 1):
            name = ws.cell(r, idx_name + 1).value
            if not name or not str(name).strip():
                continue
            if _norm_id(name) in _JUNK_IDS:
                continue
            total = ws.cell(r, idx_total + 1).value
            if not _isnum(total):
                continue
            total = float(total)
            dt = _parse_dt(ws.cell(r, idx_date + 1).value) if idx_date is not None else None
            if dt:
                dates.append(dt)
            subjects = {}
            for qi, c in enumerate(q_cols, 1):
                v = ws.cell(r, c).value
                if _isnum(v):
                    subjects[f"Q{qi}"] = float(v)
            rid = _norm_id(name)
            rows.append({"rep_id": rid, "name": str(name).strip(),
                         "total": total, "subjects": subjects})
            max_total = max(max_total, total)

        # 同人同场去重：保留最高分
        best = {}
        for row in rows:
            k = row["rep_id"]
            if k not in best or row["total"] > best[k]["total"]:
                best[k] = row
        rows = list(best.values())
        if not rows:
            raise ValueError("未解析到任何有效成绩行")

        # 检测「总分」列是否为比例值(0~1)。若是且填了「固定满分」，则换算为绝对分；
        # 否则按原值(绝对分)处理。避免比例值被当作绝对分直接比较导致全员不通过。
        _max_total = max((row["total"] for row in rows), default=0)
        ratio_warning = ""
        if _max_total > 0 and _max_total <= 1.0:
            if full_score:
                _fs = float(full_score)
                for row in rows:
                    row["total"] = round(row["total"] * _fs, 4)
                ratio_warning = f"；检测到总分列为比例值(0-1)，已按满分 {_fs} 换算为绝对分"
            else:
                ratio_warning = ("；警告：总分列为比例值(0-1)但未填「固定满分」，将按原小数作绝对分比较，"
                                 "可能全员不通过——请填写「固定满分」后重新导入")
            _max_total = max((row["total"] for row in rows), default=0)
            max_total = _max_total

        if not exam_date:
            exam_date = min(dates).strftime("%Y-%m-%d") if dates else _now()[:10]
        if pass_score is not None:
            ps = float(pass_score)
        elif full_score is not None:
            ps = round(float(full_score) * pass_ratio, 1)
        else:
            ps = round(max_total * pass_ratio, 1) if max_total else 60

        # 建缺失客服（满足外键）
        for row in rows:
            exists = self._row(self.conn.execute(
                "SELECT rep_id FROM cs_reps WHERE rep_id=?", (row["rep_id"],)))
            if not exists:
                self.create_rep({"rep_id": row["rep_id"], "name": row["name"], "hire_date": None})

        note = f"网页导入 {os.path.basename(path)}，共 {len(q_cols)} 道小题"
        results = [{"rep_id": x["rep_id"], "name": x["name"],
                    "subjects": x["subjects"], "total": x["total"],
                    **({"full_score": full_score} if full_score is not None else {})} for x in rows]
        sess = self.create_session({
            "exam_name": exam_name, "batch": batch, "exam_date": exam_date,
            "pass_score": ps, "note": note}, results=results)
        return {"session_id": sess["session_id"], "exam_name": exam_name,
                "batch": batch, "exam_date": exam_date, "count": len(rows),
                "max_total": max_total, "pass_score": ps, "q_count": len(q_cols),
                "warning": ratio_warning}

    # ---- 知识维度 ----
    def list_dimensions(self) -> list:
        return self._rows(self.conn.execute(
            "SELECT * FROM knowledge_dimensions ORDER BY dim_id"))

    def create_dimension(self, data: dict) -> dict:
        if not data.get("name_cn"):
            raise ValueError("维度中文名必填")
        cur = self.conn.execute(
            "INSERT INTO knowledge_dimensions (name_cn, name_en, description, created_at) "
            "VALUES (?,?,?,?)",
            (data["name_cn"].strip(), data.get("name_en"), data.get("description"), _now()))
        self.conn.commit()
        return self._row(self.conn.execute(
            "SELECT * FROM knowledge_dimensions WHERE dim_id=?", (cur.lastrowid,)))

    def update_dimension(self, dim_id: int, data: dict) -> dict:
        fields, vals = [], []
        for k in ("name_cn", "name_en", "description"):
            if k in data:
                v = data[k]
                if isinstance(v, str):
                    v = v.strip()
                fields.append(f"{k}=?"); vals.append(v)
        if fields:
            vals.append(dim_id)
            self.conn.execute(
                f"UPDATE knowledge_dimensions SET {','.join(fields)} WHERE dim_id=?",
                vals)
            self.conn.commit()
        return self._row(self.conn.execute(
            "SELECT * FROM knowledge_dimensions WHERE dim_id=?", (dim_id,))) or {}

    def delete_dimension(self, dim_id: int) -> None:
        # 删除维度时，连带清理其题→维度映射
        self.conn.execute("DELETE FROM exam_question_dimensions WHERE dim_id=?", (dim_id,))
        self.conn.execute("DELETE FROM knowledge_dimensions WHERE dim_id=?", (dim_id,))
        self.conn.commit()

    def bulk_delete_dimensions(self, ids: list) -> int:
        """批量删除知识维度（每个维度连带清理其题→维度映射）。返回成功删除数。"""
        n = 0
        for did in (ids or []):
            try:
                self.delete_dimension(int(did))
                n += 1
            except Exception:
                pass
        return n

    def set_question_dimension(self, exam_name, q_number, dim_id, max_score=None) -> list:
        self.conn.execute(
            "INSERT OR REPLACE INTO exam_question_dimensions (exam_name, q_number, dim_id, max_score) "
            "VALUES (?,?,?,?)",
            (exam_name, str(q_number), int(dim_id), max_score))
        self.conn.commit()
        return self._rows(self.conn.execute(
            "SELECT * FROM exam_question_dimensions WHERE exam_name=? AND q_number=?",
            (exam_name, str(q_number))))

    def delete_question_dimension(self, exam_name, q_number, dim_id) -> None:
        self.conn.execute(
            "DELETE FROM exam_question_dimensions WHERE exam_name=? AND q_number=? AND dim_id=?",
            (exam_name, str(q_number), int(dim_id)))
        self.conn.commit()

    def get_exam_question_dimensions(self, exam_name) -> list:
        return self._rows(self.conn.execute(
            "SELECT q.*, d.name_cn, d.name_en FROM exam_question_dimensions q "
            "JOIN knowledge_dimensions d ON q.dim_id=d.dim_id "
            "WHERE q.exam_name=? ORDER BY q.q_number", (exam_name,)))

    # ---- 维度分析 ----
    def exam_dimension_distribution(self, exam_name) -> dict:
        """一场考试中各知识维度的题目占比（基于题→维度映射）。"""
        rows = self._rows(self.conn.execute(
            "SELECT q.dim_id, d.name_cn, d.name_en, COUNT(*) c "
            "FROM exam_question_dimensions q "
            "JOIN knowledge_dimensions d ON q.dim_id=d.dim_id "
            "WHERE q.exam_name=? GROUP BY q.dim_id", (exam_name,)))
        total = sum(r["c"] for r in rows) or 1
        for r in rows:
            r["ratio"] = round(r["c"] / total * 100, 1)
        return {"exam_name": exam_name, "total_mapped": total, "dimensions": rows}

    def rep_dimension_weakness(self, rep_id, session_ids=None) -> dict:
        """
        分析某客服在各知识维度的错题分布，并给出文字学习建议。
        支持按考试范围过滤：session_ids 为 None → 该客服全部考试（综合全部）；
        为列表 → 仅统计指定考试（单场错题 / 多场综合维度统计）。
        「弱题」判定：该题得分 < 该题满分估计 × 0.6。
        满分估计优先用题目设定满分(questions.score)，其次用所选范围内全员最高分，
        使「单场单人」也能识别弱题（不再依赖同场有其他考生）。
        仅统计已映射知识维度的题；未映射的题跳过。
        """
        # 考试范围过滤条件
        scope_filter = ""
        params = [rep_id]
        if session_ids:
            placeholders = ",".join("?" * len(session_ids))
            scope_filter = f" AND r.session_id IN ({placeholders})"
            params = [rep_id] + list(session_ids)

        my = self._rows(self.conn.execute(
            "SELECT r.subjects, r.session_id, s.exam_name, s.exam_date FROM exam_results r "
            f"JOIN exam_sessions s ON r.session_id=s.session_id WHERE r.rep_id=?{scope_filter}",
            params))
        if not my:
            return {}
        # 预计算每 (exam_name, q_number) 的全员最高分，限定在所选考试范围内
        qmax = {}
        all_params = list(session_ids) if session_ids else []
        all_filter = f" AND r.session_id IN ({','.join('?' * len(session_ids))})" if session_ids else ""
        allrows = self._rows(self.conn.execute(
            "SELECT r.subjects, s.exam_name FROM exam_results r "
            f"JOIN exam_sessions s ON r.session_id=s.session_id WHERE 1=1{all_filter}",
            all_params))
        for r in allrows:
            subj = json.loads(r["subjects"]) if r["subjects"] else {}
            for q, v in subj.items():
                try:
                    v = float(v)
                except (TypeError, ValueError):
                    continue
                key = (r["exam_name"], str(q))
                qmax[key] = max(qmax.get(key, 0), v)

        # 若题→维度映射里记录了该题满分(max_score)，优先用它作为满分估计
        for m in self._rows(self.conn.execute(
                "SELECT exam_name, q_number, max_score FROM exam_question_dimensions "
                "WHERE max_score IS NOT NULL")):
            key = (m["exam_name"], str(m["q_number"]))
            if m["max_score"]:
                qmax[key] = max(qmax.get(key, 0), float(m["max_score"]))

        # 预构建 题→维度名 映射，供弱题展示其全部归属维度（自由组合）
        qdims_map = {}
        for m in self._rows(self.conn.execute(
                "SELECT eq.exam_name, eq.q_number, k.name_cn FROM exam_question_dimensions eq "
                "JOIN knowledge_dimensions k ON eq.dim_id=k.dim_id")):
            qdims_map.setdefault((m["exam_name"], str(m["q_number"])), []).append(m["name_cn"])

        dim_hits = {}
        total_weak = 0
        for r in my:
            subj = json.loads(r["subjects"]) if r["subjects"] else {}
            ename = r["exam_name"]
            for q, v in subj.items():
                try:
                    v = float(v)
                except (TypeError, ValueError):
                    continue
                qd = self.get_question_by_exam_q(ename, str(q))
                # 满分估计：优先用题目设定分（PPT 导入已写 score；单条/单人考试也成立），
                # 仅在题目分缺失时回退到所选范围的全员最高分。
                # 满分估计必须与 subjects 实际分数尺度一致，否则会出现 1/5 这类错配、
                # 并把答对的题误判为弱题（弱项列表被夸大 → 看起来像“全部显示”）。
                # 优先级：所选范围全员实测最高分(qmax, 天然同尺度) → 题库设定分(兜底)。
                # 满分估计优先用题库真实分（questions.score，即用户设定分），
                # 其次用所选范围全员实测最高分（qmax，天然同尺度）。
                # 不再优先 exam_question_dimensions.max_score（在线考试交卷时会
                # 被写成组卷默认分5，导致错题分母误显示为5）。
                full = float(qd["score"]) if (qd and qd.get("score")) else None
                if not full or full <= 0:
                    full = qmax.get((ename, str(q)))
                if not full or full <= 0:
                    continue
                if v < full * 0.6:
                    maps = self._rows(self.conn.execute(
                        "SELECT dim_id FROM exam_question_dimensions WHERE exam_name=? AND q_number=?",
                        (ename, str(q))))
                    if not maps:
                        continue
                    for m in maps:
                        did = m["dim_id"]
                        dim = self._row(self.conn.execute(
                            "SELECT name_cn, name_en FROM knowledge_dimensions WHERE dim_id=?", (did,)))
                        d = dim_hits.setdefault(did, {
                            "dim_id": did, "name_cn": dim["name_cn"], "name_en": dim["name_en"],
                            "weak_count": 0, "questions": []})
                        d["weak_count"] += 1
                        rec = {"exam_name": ename, "q": str(q), "score": v, "max": full,
                               "dims": qdims_map.get((ename, str(q)), [])}
                        # 关联 PPT 题库：补全题干/正确答案/解析，供学习建议展示
                        if qd:
                            rec["content"] = qd["content"]
                            rec["answer"] = qd["answer"]
                            rec["q_type"] = qd["q_type"]
                            try:
                                rec["options"] = normalize_options(qd["options"])
                            except (TypeError, ValueError):
                                rec["options"] = None
                            rec["explanation"] = qd["explanation"]
                        d["questions"].append(rec)
                        total_weak += 1

        ranking = sorted(dim_hits.values(), key=lambda x: -x["weak_count"])

        # 本次分析覆盖的考试范围
        # 注意：summary 文字不再在服务端生成（避免混入中文）。前端按当前 UI 语言
        # 用 ranking（含 name_cn / name_en 与 questions）动态拼装，保证中英一致。
        seen = {}
        for r in my:
            seen[r["session_id"]] = {"session_id": r["session_id"],
                                     "exam_name": r["exam_name"], "exam_date": r["exam_date"]}
        exams = sorted(seen.values(), key=lambda x: (x["exam_date"] or ""))
        if session_ids and len(session_ids) == 1:
            scope = "single"
        elif session_ids and len(session_ids) > 1:
            scope = "combined"
        else:
            scope = "all"
        return {
            "rep_id": rep_id,
            "scope": scope,
            "exams": exams,
            "total_weak": total_weak,
            "ranking": ranking,
        }

    # ---- 看板聚合 ----
    def individual_view(self, rep_id: str) -> dict:
        rep = self._row(self.conn.execute(
            "SELECT * FROM cs_reps WHERE rep_id=?", (rep_id,)))
        if not rep:
            return {}
        rows = self._rows(self.conn.execute(
            "SELECT r.*, s.exam_name, s.batch, s.exam_date, s.pass_score, s.exam_type "
            "FROM exam_results r JOIN exam_sessions s ON r.session_id=s.session_id "
            "WHERE r.rep_id=? ORDER BY s.exam_date", (rep_id,)))
        for row in rows:
            row["subjects"] = json.loads(row["subjects"]) if row["subjects"] else {}
        totals = [r["total"] for r in rows if r["total"] is not None]
        trend = [{"exam_date": r["exam_date"], "exam_name": r["exam_name"],
                  "batch": r["batch"], "exam_type": r["exam_type"], "total": r["total"]} for r in rows]
        delta = (totals[-1] - totals[0]) if len(totals) >= 2 else None
        span = None
        if rep.get("hire_date") and rows:
            span = {"hire_date": rep.get("hire_date"),
                    "leave_date": rep.get("leave_date"),
                    "status": rep.get("status"),
                    "first_exam": rows[0]["exam_date"],
                    "latest_exam": rows[-1]["exam_date"]}
        return {
            "rep": rep,
            "records": rows,
            "trend": trend,
            "summary": {
                "first_exam": rows[0]["exam_date"] if rows else None,
                "latest_exam": rows[-1]["exam_date"] if rows else None,
                "exam_count": len(rows),
                "first_total": totals[0] if totals else None,
                "latest_total": totals[-1] if totals else None,
                "delta": round(delta, 2) if delta is not None else None,
                "pass_count": sum(1 for r in rows if r["passed"]),
            },
            "lifecycle": span,
        }

    def batch_view(self, session_id: int) -> dict:
        session = self.get_session(session_id)
        if not session:
            return {}
        rows = self._rows(self.conn.execute(
            "SELECT r.*, s.exam_name, s.batch, s.exam_date, s.pass_score, s.exam_type "
            "FROM exam_results r JOIN exam_sessions s ON r.session_id=s.session_id "
            "WHERE r.session_id=? ORDER BY r.total DESC", (session_id,)))
        for row in rows:
            row["subjects"] = json.loads(row["subjects"]) if row["subjects"] else {}
        # 成绩分布：按得分率(%)分桶，与通过判定口径一致（在线考试满分可能仅10分，
        # 用绝对总分分桶会把 9/10 也误判进 <60）
        buckets = {"<60%": 0, "60-69%": 0, "70-79%": 0, "80-89%": 0, "90-100%": 0}
        for r in rows:
            pct = self._rate_pct(r)
            if pct is None:
                continue
            if pct < 60: buckets["<60%"] += 1
            elif pct < 70: buckets["60-69%"] += 1
            elif pct < 80: buckets["70-79%"] += 1
            elif pct < 90: buckets["80-89%"] += 1
            else: buckets["90-100%"] += 1
        # 批次维度占比（基于题→维度映射）
        dim_dist = self.exam_dimension_distribution(session["exam_name"])
        return {
            "session": session,
            "records": rows,
            "distribution": buckets,
            "stats": self._session_stats(session_id),
            "dimension_distribution": dim_dist,
            "no_shows": self.batch_no_shows(session_id),
        }

    def _rate_pct(self, row: dict):
        """把一行成绩换算成得分率百分比（0-100）。优先用 score_rate，否则 total/full_score。"""
        sr = row.get("score_rate")
        if sr is not None:
            return sr * 100 if sr <= 1.0 else sr
        fs = row.get("full_score")
        if fs:
            t = row.get("total")
            if t is not None:
                return (t / fs) * 100
        return None

    def batch_no_shows(self, session_id: int) -> list:
        """返回「被分配本场考试但未参加」的客服名单（仅在线考试适用）。

        通过 session.note 形如 'online:paper_id=N' 解析出试卷，再比对
        exam_assignments（应考）与 exam_attempts（实考）。Excel/手动导入批次
        没有 paper 关联，无法判定，返回空列表。
        """
        session = self.get_session(session_id)
        if not session:
            return []
        note = session.get("note") or ""
        m = re.search(r"paper_id=(\d+)", note)
        if not m:
            return []
        paper_id = int(m.group(1))
        assigned = self.get_paper_assignments(paper_id)
        if not assigned:
            return []  # 未指定分配=全员可见，无法判定缺席
        taken = {r["rep_id"] for r in self._rows(self.conn.execute(
            "SELECT DISTINCT rep_id FROM exam_attempts WHERE paper_id=?", (paper_id,)))}
        missing = [rid for rid in assigned if rid not in taken]
        if not missing:
            return []
        reps = {r["rep_id"]: r["name"] for r in self._rows(self.conn.execute(
            "SELECT rep_id, name FROM cs_reps WHERE rep_id IN (%s)" % ",".join("?" * len(missing)),
            missing))}
        return [{"rep_id": rid, "name": reps.get(rid, "")} for rid in missing]

    def period_view(self, start: str, end: str) -> dict:
        rows = self._rows(self.conn.execute(
            "SELECT r.total, r.passed, s.exam_date, s.exam_name, s.batch, s.session_id, s.exam_type "
            "FROM exam_results r JOIN exam_sessions s ON r.session_id=s.session_id "
            "WHERE s.exam_date>=? AND s.exam_date<=? ORDER BY s.exam_date",
            (start, end)))
        # 按场次聚合
        by_session = {}
        for r in rows:
            sid = r["session_id"]
            by_session.setdefault(sid, {
                "session_id": sid, "exam_date": r["exam_date"],
                "exam_name": r["exam_name"], "batch": r["batch"], "exam_type": r["exam_type"],
                "totals": [], "passed": 0, "count": 0})
            by_session[sid]["totals"].append(r["total"])
            by_session[sid]["passed"] += r["passed"]
            by_session[sid]["count"] += 1
        series = []
        for sid, d in by_session.items():
            series.append({
                "session_id": sid, "exam_date": d["exam_date"],
                "exam_name": d["exam_name"], "batch": d["batch"], "exam_type": d["exam_type"],
                "avg_total": round(sum(d["totals"]) / len(d["totals"]), 2),
                "pass_rate": round(d["passed"] / d["count"] * 100, 1),
                "count": d["count"],
            })
        series.sort(key=lambda x: x["exam_date"])
        return {"range": {"start": start, "end": end}, "series": series,
                "total_records": len(rows)}

    def overview(self) -> dict:
        rep_count = self._row(self.conn.execute("SELECT COUNT(*) c FROM cs_reps"))["c"]
        session_count = self._row(self.conn.execute("SELECT COUNT(*) c FROM exam_sessions"))["c"]
        result_count = self._row(self.conn.execute("SELECT COUNT(*) c, SUM(passed) p FROM exam_results"))
        avg = self._row(self.conn.execute("SELECT AVG(total) a FROM exam_results"))["a"]
        avg_sr = self._row(self.conn.execute(
            "SELECT AVG(score_rate) a FROM exam_results WHERE score_rate IS NOT NULL"))["a"]
        dim_count = self._row(self.conn.execute("SELECT COUNT(*) c FROM knowledge_dimensions"))["c"]
        mapped = self._row(self.conn.execute("SELECT COUNT(DISTINCT exam_name) c FROM exam_question_dimensions"))["c"]
        return {
            "rep_count": rep_count,
            "session_count": session_count,
            "result_count": result_count["c"],
            "pass_count": result_count["p"] or 0,
            "pass_rate": round(result_count["p"] / result_count["c"] * 100, 1) if result_count["c"] else None,
            "avg_total": round(avg, 2) if avg is not None else None,
            "avg_score_rate": round(avg_sr, 4) if avg_sr is not None else None,
            "pass_line_ratio": self.get_pass_line_ratio(),
            "dimension_count": dim_count,
            "exams_with_dimension_mapping": mapped,
        }

    def recompute_score_rates(self) -> dict:
        """历史成绩重算：补算 score_rate + 按得分率比例重标 passed。
        返回 {attempts, results, updated} 计数。"""
        ratio = self.get_pass_line_ratio()
        upd_a = 0
        cur = self.conn.execute(
            "SELECT attempt_id, paper_id, total_score FROM exam_attempts")
        for a in cur.fetchall():
            fs = self._paper_full_score(a["paper_id"])
            # 在线考试按该批次(paper 对应 session)的 pass_score 判定，而非全局比例，
            # 避免个人视图(读 exam_attempts)与管理端(读 exam_results)判定不一致。
            sess = self._row(self.conn.execute(
                "SELECT * FROM exam_sessions WHERE note=?", (f"online:paper_id={a['paper_id']}",)))
            if sess:
                ps = sess["pass_score"]
                online = True
            else:
                ps = None
                online = False
            sr, passed = self._score_rate_pass(a["total_score"], fs, ps, ratio, online=online)
            self.conn.execute(
                "UPDATE exam_attempts SET score_rate=?, passed=?, full_score=? WHERE attempt_id=?",
                (sr, passed, fs, a["attempt_id"]))
            upd_a += 1
        upd_r = 0
        cur = self.conn.execute(
            "SELECT result_id, session_id, total, full_score FROM exam_results")
        for r in cur.fetchall():
            session = self.get_session(int(r["session_id"]))
            ps = session["pass_score"] if session else 60
            note = (session.get("note") or "") if session else ""
            if note.startswith("online:"):
                # 在线考试：满分=试卷各小题分值之和（total 即其和，量纲一致）
                m = re.search(r"paper_id=(\d+)", note)
                pid = int(m.group(1)) if m else None
                fs = r["full_score"] or (self._paper_full_score(pid) if pid else None)
            else:
                # Excel/手动导入：保留导入时显式提供的满分（仅用于得分率展示），
                # 不回退题库总和（量纲不一致会算成极小得分率），也不清空已有满分。
                # 绝对分判定与满分无关。
                fs = r["full_score"]
            sr, passed = self._score_rate_pass(r["total"], fs, ps, ratio,
                                               online=note.startswith("online:"))
            self.conn.execute(
                "UPDATE exam_results SET score_rate=?, passed=?, full_score=? WHERE result_id=?",
                (sr, passed, fs, r["result_id"]))
            upd_r += 1
        self.conn.commit()
        return {"attempts": upd_a, "results": upd_r, "updated": upd_a + upd_r}

    def rep_score_rates(self, scope="avg", session_id=None) -> list:
        """各客服得分率。scope=avg 取每人平均；scope=session 按场次。"""
        ratio = self.get_pass_line_ratio()
        where = "WHERE r.score_rate IS NOT NULL"
        if scope == "session" and session_id:
            where += f" AND r.session_id={int(session_id)}"
        sql = (
            "SELECT r.rep_id, rp.name, rp.position, rp.channel, "
            "AVG(r.score_rate) avg_rate, COUNT(*) cnt, SUM(r.passed) passed_cnt "
            f"FROM exam_results r JOIN cs_reps rp ON rp.rep_id=r.rep_id {where} "
            "GROUP BY r.rep_id, rp.name, rp.position, rp.channel "
            "ORDER BY avg_rate DESC")
        rows = self._rows(self.conn.execute(sql))
        for row in rows:
            row["avg_rate"] = round(row["avg_rate"], 4) if row["avg_rate"] is not None else None
            row["pass_line_ratio"] = ratio
            row["meets"] = bool(row["avg_rate"] is not None and row["avg_rate"] >= ratio)
        return rows

    # ---- 增强3-需求2：积分系统 ----
    def get_points_rules(self) -> dict:
        try:
            d = json.loads(self.get_config("points_rules", "{}") or "{}")
        except (TypeError, ValueError):
            d = {}
        d.setdefault("participate", 10)
        d.setdefault("pass", {0.88: 20, 0.90: 30, 0.95: 40})
        d.setdefault("material", 5)
        d.setdefault("mini_quiz", 10)
        return d

    def get_points_threshold(self) -> int:
        try:
            return int(float(self.get_config("points_threshold", "100") or "100"))
        except (TypeError, ValueError):
            return 100

    def get_points_period(self) -> tuple:
        """周期目标：period=month|quarter，target=每周期需获取积分(0=未启用)。"""
        period = (self.get_config("points_period", "quarter") or "quarter")
        if period not in ("month", "quarter"):
            period = "quarter"
        try:
            target = int(float(self.get_config("points_period_target", "0") or 0))
        except (TypeError, ValueError):
            target = 0
        return period, target

    def rep_period_points(self, rep_id, period, year=None) -> int:
        """当前周期(月/季)已获积分。"""
        now = datetime.datetime.now()
        if year is None:
            year = now.year
        if period == "month":
            m = now.month
            row = self._row(self.conn.execute(
                "SELECT COALESCE(SUM(delta),0) s FROM points_log "
                "WHERE rep_id=? AND year=? AND CAST(strftime('%m', created_at) AS INTEGER)=?",
                (rep_id, year, m)))
        else:
            q = (now.month - 1) // 3 + 1
            row = self._row(self.conn.execute(
                "SELECT COALESCE(SUM(delta),0) s FROM points_log "
                "WHERE rep_id=? AND year=? AND quarter=?",
                (rep_id, year, q)))
        return (row or {}).get("s", 0)

    def pass_points_for(self, score_rate, passed=False) -> int:
        """通过考试的积分：有得分率按分层取最高档；无得分率(Excel 导入)按绝对通过给最低档。"""
        rules = self.get_points_rules()
        tiers = rules.get("pass", {})
        if not tiers:
            return 0
        if score_rate is not None:
            best = 0
            for thr, pts in tiers.items():
                if float(score_rate) >= float(thr):
                    best = max(best, int(pts))
            return best
        if passed:
            return min(int(v) for v in tiers.values())
        return 0

    def is_dim_points_claimed(self, rep_id, dim_id, year=None, quarter=None):
        """检查客服是否已在某维度的本月/本季度小测通过积分。用于防刷分。返回 True 表示已领过（不重复发）。

        当前实现：按季度去重（更简单且符合周期目标设定）。
        """
        rep_id = (rep_id or "").strip()
        if not rep_id or not dim_id:
            return True
        # 如果未指定，使用当前时间
        if year is None or quarter is None:
            from datetime import datetime
            dt = datetime.now()
            year, quarter = dt.year, (dt.month - 1) // 3 + 1
        exist = self._row(self.conn.execute(
            "SELECT 1 FROM points_log WHERE rep_id=? AND rule_key='mini_quiz' AND ref_type='quiz' AND ref_id LIKE ? AND year=? AND quarter=?",
            (rep_id, f"{dim_id}%", year, quarter)))
        return bool(exist)

    def award_points(self, rep_id, rule_key, ref_type, ref_id, note="", delta=None):
        """发放积分；同一 (ref_type, ref_id) 不重复发放。返回实际发放增量(0 表示已发过/无分)。"""
        rep_id = (rep_id or "").strip()
        if not rep_id:
            return 0
        exist = self._row(self.conn.execute(
            "SELECT 1 FROM points_log WHERE rep_id=? AND rule_key=? AND ref_type=? AND ref_id=?",
            (rep_id, rule_key, ref_type, str(ref_id))))
        if exist:
            return 0
        rules = self.get_points_rules()
        if delta is None:
            delta = int(rules.get(rule_key, 0) or 0)
        delta = int(delta or 0)
        if delta == 0:
            return 0
        self.conn.execute(
            "INSERT OR IGNORE INTO points_account (rep_id, total, updated_at) VALUES (?,0,?)",
            (rep_id, _now()))
        self.conn.execute(
            "UPDATE points_account SET total=total+?, updated_at=? WHERE rep_id=?",
            (delta, _now(), rep_id))
        now = _now()
        y = int(now[:4]); m = int(now[5:7]); q = (m - 1) // 3 + 1
        self.conn.execute(
            "INSERT INTO points_log (rep_id, rule_key, delta, ref_type, ref_id, note, created_at, year, quarter) "
            "VALUES (?,?,?,?,?,?,?,?,?)",
            (rep_id, rule_key, delta, ref_type, str(ref_id), note, now, y, q))
        self.conn.commit()
        return delta

    def points_summary(self) -> list:
        """管理端积分总览（全量历史，兼容旧调用）：每人总分 + 是否达标。"""
        threshold = self.get_points_threshold()
        rows = self._rows(self.conn.execute(
            "SELECT a.rep_id, r.name, r.position, r.channel, a.total, "
            "(SELECT COUNT(*) FROM points_log l WHERE l.rep_id=a.rep_id) cnt "
            "FROM points_account a LEFT JOIN cs_reps r ON r.rep_id=a.rep_id "
            "ORDER BY a.total DESC, a.rep_id"))
        for row in rows:
            row["threshold"] = threshold
            row["meets"] = bool(row["total"] is not None and row["total"] >= threshold)
        return rows

    def points_year_summary(self, year: int = None) -> list:
        """管理端积分总览（按年度）：每人 Q1~Q4 各季度积分 + 年度合计 + 是否达标。
        year 缺省取当前年。跨年自动归零（无记录即 0）；历史年可筛选查看。"""
        if year is None:
            year = datetime.datetime.now().year
        year = int(year)
        threshold = self.get_points_threshold()
        rows = self._rows(self.conn.execute(
            "SELECT r.rep_id, r.name, r.position, r.channel, "
            "COALESCE(SUM(CASE WHEN l.quarter=1 THEN l.delta ELSE 0 END),0) q1, "
            "COALESCE(SUM(CASE WHEN l.quarter=2 THEN l.delta ELSE 0 END),0) q2, "
            "COALESCE(SUM(CASE WHEN l.quarter=3 THEN l.delta ELSE 0 END),0) q3, "
            "COALESCE(SUM(CASE WHEN l.quarter=4 THEN l.delta ELSE 0 END),0) q4, "
            "COALESCE(SUM(l.delta),0) year_total "
            "FROM cs_reps r LEFT JOIN points_log l ON l.rep_id=r.rep_id AND l.year=? "
            "GROUP BY r.rep_id, r.name, r.position, r.channel "
            "ORDER BY year_total DESC, r.rep_id", (year,)))
        for row in rows:
            row["threshold"] = threshold
            row["meets"] = bool(row["year_total"] >= threshold)
            row["year"] = year
            pperiod, ptarget = self.get_points_period()
            row["period"] = pperiod
            row["period_target"] = ptarget
            row["period_points"] = self.rep_period_points(row["rep_id"], pperiod, year) if ptarget > 0 else 0
            row["period_meets"] = bool(ptarget > 0 and row["period_points"] >= ptarget)
        return rows

    def list_point_years(self) -> list:
        """可筛选的年份列表：日志中出现过的年份 + 当前年（降序）。"""
        cur_year = datetime.datetime.now().year
        yrs = [r[0] for r in self._rows(self.conn.execute(
            "SELECT DISTINCT year FROM points_log WHERE year IS NOT NULL"))]
        if cur_year not in yrs:
            yrs.append(cur_year)
        return sorted(set(yrs), reverse=True)

    def rep_points(self, rep_id, year=None) -> dict:
        if year is not None:
            year = int(year)
            row = self._row(self.conn.execute(
                "SELECT COALESCE(SUM(CASE WHEN quarter=1 THEN delta ELSE 0 END),0) q1, "
                "COALESCE(SUM(CASE WHEN quarter=2 THEN delta ELSE 0 END),0) q2, "
                "COALESCE(SUM(CASE WHEN quarter=3 THEN delta ELSE 0 END),0) q3, "
                "COALESCE(SUM(CASE WHEN quarter=4 THEN delta ELSE 0 END),0) q4, "
                "COALESCE(SUM(delta),0) total FROM points_log WHERE rep_id=? AND year=?",
                (rep_id, year)))
            threshold = self.get_points_threshold()
            total = (row or {}).get("total", 0)
            pperiod, ptarget = self.get_points_period()
            return {"rep_id": rep_id, "year": year, "q1": row["q1"], "q2": row["q2"],
                    "q3": row["q3"], "q4": row["q4"], "total": total,
                    "threshold": threshold, "meets": bool(total >= threshold),
                    "period": pperiod, "period_target": ptarget,
                    "period_points": self.rep_period_points(rep_id, pperiod, year) if ptarget > 0 else 0,
                    "period_meets": bool(ptarget > 0 and self.rep_period_points(rep_id, pperiod, year) >= ptarget)}
        row = self._row(self.conn.execute(
            "SELECT * FROM points_account WHERE rep_id=?", (rep_id,)))
        pperiod, ptarget = self.get_points_period()
        if not row:
            return {"rep_id": rep_id, "total": 0, "meets": False, "threshold": self.get_points_threshold(),
                    "period": pperiod, "period_target": ptarget, "period_points": 0, "period_meets": False}
        row["threshold"] = self.get_points_threshold()
        row["meets"] = bool(row["total"] >= row["threshold"])
        row["period"] = pperiod
        row["period_target"] = ptarget
        row["period_points"] = self.rep_period_points(rep_id, pperiod) if ptarget > 0 else 0
        row["period_meets"] = bool(ptarget > 0 and row["period_points"] >= ptarget)
        return row

    def rep_points_log(self, rep_id, year=None) -> list:
        if year is not None:
            return self._rows(self.conn.execute(
                "SELECT * FROM points_log WHERE rep_id=? AND year=? "
                "ORDER BY created_at DESC, log_id DESC", (rep_id, int(year))))
        return self._rows(self.conn.execute(
            "SELECT * FROM points_log WHERE rep_id=? ORDER BY created_at DESC, log_id DESC",
            (rep_id,)))

    def update_points_config(self, rules: dict = None, threshold: int = None,
                              period: str = None, period_target: int = None) -> dict:
        """保存积分规则 + 阈值 + 周期目标（系统设置视图调用）。"""
        if rules is not None:
            # 规整 pass 分层为 float->int
            norm = {"participate": int(rules.get("participate", 10)),
                    "material": int(rules.get("material", 5)),
                    "mini_quiz": int(rules.get("mini_quiz", 10))}
            pass_tiers = {}
            for k, v in (rules.get("pass") or {}).items():
                try:
                    pass_tiers[float(k)] = int(v)
                except (TypeError, ValueError):
                    continue
            norm["pass"] = pass_tiers or {0.88: 20, 0.90: 30, 0.95: 40}
            self.set_config("points_rules", json.dumps(norm, ensure_ascii=False))
        if threshold is not None:
            try:
                self.set_config("points_threshold", int(threshold))
            except (TypeError, ValueError):
                pass  # 上游 api 已做校验，这里仅兜底
        if period is not None:
            self.set_config("points_period", period if period in ("month", "quarter") else "quarter")
        if period_target is not None:
            try:
                self.set_config("points_period_target", int(period_target))
            except (TypeError, ValueError):
                pass
        return {"rules": self.get_points_rules(), "threshold": self.get_points_threshold(),
                "period": self.get_points_period()[0], "period_target": self.get_points_period()[1]}

    @staticmethod
    def _parse_ans(a):
        if a is None:
            return []
        if isinstance(a, (list, tuple)):
            a = ",".join(str(x) for x in a)
        return [x.strip().upper() for x in str(a).replace("；", ",").replace(";", ",").split(",") if x.strip()]

    # ---- 学习资料库 / 智能推荐 / 小测 ----
    @staticmethod
    def _parse_dim_ids(raw):
        """兼容多种入参：列表 / 逗号串 / 单数字，返回去重整数列表。"""
        if raw is None:
            return []
        if isinstance(raw, str):
            raw = raw.split(",")
        out = []
        for x in raw:
            try:
                x = int(str(x).strip())
            except (TypeError, ValueError):
                continue
            if x and x not in out:
                out.append(x)
        return out

    def create_material(self, data: dict) -> dict:
        dim_ids = self._parse_dim_ids(data.get("dim_ids") or data.get("dim_id"))
        # 兼容旧逻辑：单维度写入 learning_materials.dim_id，其余走关联表
        legacy = dim_ids[0] if dim_ids else None
        cur = self.conn.execute(
            "INSERT INTO learning_materials (title, mtype, dim_id, content, file_path, url, link_kind, created_at) "
            "VALUES (?,?,?,?,?,?,?,?)",
            (data.get("title", "").strip(), data.get("mtype", "text"),
             legacy, data.get("content"), data.get("file_path"),
             data.get("url"), data.get("link_kind") or None, _now()))
        mid = cur.lastrowid
        for did in dim_ids:
            self.conn.execute(
                "INSERT OR IGNORE INTO material_dimensions (material_id, dim_id) VALUES (?,?)",
                (mid, did))
        self.conn.commit()
        return self.get_material(mid)

    def update_material(self, mid, data: dict) -> dict:
        # 基础字段更新（标题/类型/内容/链接/链接类型，可选文件）
        fields, vals = [], []
        for col in ("title", "mtype", "content", "url", "link_kind"):
            if col in data and data[col] is not None:
                vals.append(data[col] if col != "title" else (data[col] or "").strip())
                fields.append(f"{col}=?")
        if "file_path" in data and data["file_path"]:
            fields.append("file_path=?"); vals.append(data["file_path"])
        if fields:
            vals.append(mid)
            self.conn.execute(
                f"UPDATE learning_materials SET {','.join(fields)} WHERE material_id=?",
                vals)
        # 维度：整表替换（兼容新增/删除维度）
        if "dim_ids" in data:
            dim_ids = self._parse_dim_ids(data.get("dim_ids"))
            self.conn.execute("DELETE FROM material_dimensions WHERE material_id=?", (mid,))
            for did in dim_ids:
                self.conn.execute(
                    "INSERT OR IGNORE INTO material_dimensions (material_id, dim_id) VALUES (?,?)",
                    (mid, did))
            legacy = dim_ids[0] if dim_ids else None
            self.conn.execute(
                "UPDATE learning_materials SET dim_id=? WHERE material_id=?", (legacy, mid))
        self.conn.commit()
        return self.get_material(mid)

    def list_materials(self, dim_id=None) -> list:
        if dim_id:
            base = self._rows(self.conn.execute(
                "SELECT m.* FROM learning_materials m "
                "JOIN material_dimensions md ON m.material_id=md.material_id "
                "WHERE md.dim_id=? ORDER BY m.material_id DESC", (int(dim_id),)))
        else:
            base = self._rows(self.conn.execute(
                "SELECT * FROM learning_materials ORDER BY material_id DESC"))
        # 聚合每个资料的维度列表（material_dimensions），使多维度全部可被前端展示；
        # 否则前端列表只能回退到 learning_materials.dim_id（单值），多维度会丢失。
        out = []
        for r in base:
            r = dict(r)
            drows = self._rows(self.conn.execute(
                "SELECT dim_id FROM material_dimensions WHERE material_id=? ORDER BY dim_id",
                (r["material_id"],)))
            r["dim_ids"] = [d["dim_id"] for d in drows]
            out.append(r)
        return out

    def get_material(self, mid) -> dict:
        m = self._row(self.conn.execute(
            "SELECT * FROM learning_materials WHERE material_id=?", (mid,))) or {}
        if m:
            m = dict(m)
            rows = self._rows(self.conn.execute(
                "SELECT dim_id FROM material_dimensions WHERE material_id=?", (mid,)))
            m["dim_ids"] = [r["dim_id"] for r in rows]
        return m

    def delete_material(self, mid) -> None:
        self.conn.execute("DELETE FROM learning_materials WHERE material_id=?", (mid,))
        self.conn.execute("DELETE FROM material_dimensions WHERE material_id=?", (mid,))
        self.conn.execute("DELETE FROM study_records WHERE material_id=?", (mid,))
        self.conn.commit()

    def bulk_delete_materials(self, ids: list) -> int:
        """批量删除学习资料（每个含关联维度与学习记录）。返回成功删除数。"""
        n = 0
        for mid in (ids or []):
            try:
                self.delete_material(int(mid))
                n += 1
            except Exception:
                pass
        return n

    def mark_material_opened(self, rep_id, material_id) -> bool:
        exist = self._row(self.conn.execute(
            "SELECT 1 FROM study_records WHERE rep_id=? AND material_id=?", (rep_id, material_id)))
        if not exist:
            self.conn.execute(
                "INSERT OR IGNORE INTO study_records (rep_id, material_id, status, progress, created_at) "
                "VALUES (?,?, 'opened', 100, ?)", (rep_id, material_id, _now()))
            self.conn.commit()
        return True

    def complete_material(self, rep_id, material_id) -> int:
        self.conn.execute(
            "INSERT OR REPLACE INTO study_records (rep_id, material_id, status, progress, created_at) "
            "VALUES (?,?, 'completed', 100, ?)", (rep_id, material_id, _now()))
        self.conn.commit()
        return self.award_points(rep_id, "material", "material", material_id, "学完资料")

    def draw_quiz_questions(self, dim_id, n) -> list:
        """随机抽取某维度客观题（不返回正确答案/解析，由服务端判分）。"""
        rows = self._rows(self.conn.execute(
            "SELECT q.question_id, q.q_type, q.content AS question, q.options, q.score, q.dim_id "
            "FROM questions q WHERE q.dim_id=? AND q.q_type IN ('single','multiple') "
            "ORDER BY RANDOM() LIMIT ?", (dim_id, int(n))))
        for r in rows:
            try:
                r["options"] = normalize_options(r["options"])
            except (TypeError, ValueError):
                r["options"] = []
        return rows

    def create_mini_quiz(self, rep_id, dim_id, question_ids) -> int:
        cur = self.conn.execute(
            "INSERT INTO mini_quiz (rep_id, dim_id, question_ids, passed, created_at) "
            "VALUES (?,?,?,0,?)",
            (rep_id, dim_id, json.dumps(question_ids, ensure_ascii=False), _now()))
        self.conn.commit()
        return cur.lastrowid

    def grade_mini_quiz(self, quiz_id, answers) -> dict:
        quiz = self._row(self.conn.execute("SELECT * FROM mini_quiz WHERE quiz_id=?", (quiz_id,)))
        if not quiz:
            return {"error": "quiz not found"}
        try:
            qids = json.loads(quiz["question_ids"]) if quiz["question_ids"] else []
        except (TypeError, ValueError):
            qids = []
        if not qids:
            return {"error": "quiz has no questions"}
        ph = ",".join("?" * len(qids))
        rows = self._rows(self.conn.execute(
            f"SELECT question_id, q_type, answer, score FROM questions WHERE question_id IN ({ph})", qids))
        correct = 0
        total = len(rows)
        for r in rows:
            qid = str(r["question_id"])
            correct_ans = self._parse_ans(r["answer"])
            student = self._parse_ans(answers.get(qid) if isinstance(answers, dict) else None)
            if set(correct_ans) == set(student):
                correct += 1
        score_rate = round(correct / total, 4) if total else 0
        ratio = self.get_pass_line_ratio()
        passed = bool(score_rate >= ratio)
        self.conn.execute("UPDATE mini_quiz SET passed=? WHERE quiz_id=?", (1 if passed else 0, quiz_id))
        awarded = 0
        if passed:
            # 防刷分：每维度每月仅首次通过给分
            if not self.is_dim_points_claimed(quiz["rep_id"], quiz["dim_id"]):
                awarded = self.award_points(quiz["rep_id"], "mini_quiz", "quiz", quiz_id, "小测通过")
            # 通过小测即视为学完该维度资料（按关联表查该维度全部资料）
            mats = self._rows(self.conn.execute(
                "SELECT DISTINCT m.material_id FROM learning_materials m "
                "JOIN material_dimensions md ON m.material_id=md.material_id "
                "WHERE md.dim_id=?", (quiz["dim_id"],)))
            for m in mats:
                self.complete_material(quiz["rep_id"], m["material_id"])
        self.conn.commit()
        return {"quiz_id": quiz_id, "correct": correct, "total": total,
                "score_rate": score_rate, "passed": passed, "points_awarded": awarded}

    def get_recommend_top_n(self) -> int:
        try:
            return int(float(self.get_config("recommend_top_n", "3") or "3"))
        except (TypeError, ValueError):
            return 3

    def get_recommend_quiz_n(self) -> int:
        try:
            return int(float(self.get_config("recommend_quiz_n", "5") or "5"))
        except (TypeError, ValueError):
            return 5

    def recommend_for_rep(self, rep_id, top_n=None, quiz_n=None) -> dict:
        """按考试场次分组推荐：每场考试（session）分别列出其弱项维度、对应学习资料与
        该维度小测的完成状态。数量为 None 时读系统设置。"""
        if top_n is None:
            top_n = self.get_recommend_top_n()
        if quiz_n is None:
            quiz_n = self.get_recommend_quiz_n()
        # 该客服的全部考试场次（在线 + 导入），按时间倒序
        results = self._rows(self.conn.execute(
            "SELECT r.session_id, s.exam_name, s.exam_date, r.passed "
            "FROM exam_results r JOIN exam_sessions s ON r.session_id=s.session_id "
            "WHERE r.rep_id=? ORDER BY r.created_at DESC", (rep_id,)))
        seen = set()
        sessions_out = []
        for r in results:
            sid = r["session_id"]
            if sid in seen:
                continue
            seen.add(sid)
            wk = self.rep_dimension_weakness(rep_id, session_ids=[sid]) or {}
            if not wk.get("ranking"):
                continue
            dims = sorted(wk["ranking"], key=lambda d: -d.get("weak_count", 0))[:top_n]
            dim_out = []
            seen_mat = set()
            for d in dims:
                did = d["dim_id"]
                mats = self.list_materials(did)
                uniq = []
                for m in mats:
                    if m["material_id"] in seen_mat:
                        continue
                    seen_mat.add(m["material_id"])
                    uniq.append(m)
                qz = self._row(self.conn.execute(
                    "SELECT passed FROM mini_quiz WHERE rep_id=? AND dim_id=? ORDER BY quiz_id DESC", (rep_id, did)))
                dim_out.append({
                    "dim_id": did,
                    "name_cn": d.get("name_cn"), "name_en": d.get("name_en"),
                    "weak_count": d.get("weak_count", 0),
                    "materials": uniq,
                    "quiz": {"done": bool(qz), "passed": (qz["passed"] == 1 if qz else None)},
                })
            sessions_out.append({
                "session_id": sid,
                "exam_name": r["exam_name"],
                "exam_date": r["exam_date"],
                "passed": r["passed"],
                "dims": dim_out,
            })
        return {"sessions": sessions_out, "quiz_n": quiz_n, "top_n": top_n}

    # ---- 题库 / 在线考试 ----
    # ---------- 题库 ----------
    def list_questions(self, filters=None):
        filters = filters or {}
        sql = ("SELECT q.*, d.name_cn AS dim_cn, d.name_en AS dim_en "
               "FROM questions q LEFT JOIN knowledge_dimensions d ON q.dim_id=d.dim_id WHERE 1=1")
        vals = []
        if filters.get("q_type"):
            sql += " AND q.q_type=?"; vals.append(filters["q_type"])
        if filters.get("category"):
            sql += " AND q.category=?"; vals.append(filters["category"])
        if filters.get("dim_id"):
            sql += " AND q.dim_id=?"; vals.append(int(filters["dim_id"]))
        if filters.get("keyword"):
            sql += " AND q.content LIKE ?"; vals.append(f"%{filters['keyword']}%")
        sql += " ORDER BY q.question_id DESC"
        return self._with_attachments(self._rows(self.conn.execute(sql, vals)))

    def create_question(self, data):
        cur = self.conn.execute(
            "INSERT INTO questions (q_type, category, content, options, answer, dim_id, score, explanation, source_exam, created_at) "
            "VALUES (?,?,?,?,?,?,?,?,?,?)",
            (data.get("q_type", "single"), data.get("category"),
             data["content"].strip(), data.get("options"),
             data.get("answer"), data.get("dim_id"), float(data.get("score") or 5),
             data.get("explanation"), (data.get("source_exam") or "").strip() or None, _now()))
        self.conn.commit()
        return self._row(self.conn.execute("SELECT * FROM questions WHERE question_id=?", (cur.lastrowid,)))

    def update_question(self, qid, data):
        fields, vals = [], []
        for k in ("q_type", "category", "content", "options", "answer", "dim_id", "score", "explanation", "source_exam"):
            if k in data:
                v = data[k]
                if k == "content" and isinstance(v, str):
                    v = v.strip()
                if k == "source_exam":
                    v = (v or "").strip() or None
                if k == "score":
                    v = float(v or 5)
                fields.append(f"{k}=?"); vals.append(v)
        if fields:
            vals.append(qid)
            self.conn.execute(f"UPDATE questions SET {','.join(fields)} WHERE question_id=?", vals)
            self.conn.commit()
        return self._row(self.conn.execute("SELECT * FROM questions WHERE question_id=?", (qid,))) or {}

    def delete_question(self, qid):
        self.conn.execute("DELETE FROM paper_questions WHERE question_id=?", (qid,))
        self.conn.execute("DELETE FROM question_attachments WHERE question_id=?", (qid,))
        self.conn.execute("DELETE FROM questions WHERE question_id=?", (qid,))
        self.conn.commit()

    def bulk_set_question_source(self, ids, source_exam):
        """批量给题目写入同一来源（source_exam）。ids 为空返回 0。"""
        ids = [int(x) for x in (ids or []) if str(x).strip()]
        if not ids:
            return 0
        val = (source_exam or "").strip() or None
        placeholders = ",".join("?" * len(ids))
        self.conn.execute(
            f"UPDATE questions SET source_exam=? WHERE question_id IN ({placeholders})",
            [val] + ids)
        self.conn.commit()
        return len(ids)

    def bulk_delete_questions(self, ids):
        """批量删除题目（含关联 paper_questions / question_attachments 及物理文件）。
        返回实际删除的题目数；ids 为空返回 0。"""
        ids = [int(x) for x in (ids or []) if str(x).strip()]
        if not ids:
            return 0
        placeholders = ",".join("?" * len(ids))
        # 附件已存数据库，随 question_attachments 行一并删除，无需物理删文件
        self.conn.execute(f"DELETE FROM paper_questions WHERE question_id IN ({placeholders})", ids)
        self.conn.execute(f"DELETE FROM question_attachments WHERE question_id IN ({placeholders})", ids)
        cur = self.conn.execute(f"DELETE FROM questions WHERE question_id IN ({placeholders})", ids)
        self.conn.commit()
        return cur.rowcount

    # ---------- 题目附件（图片等） ----------
    def add_question_attachment(self, qid, file_obj, filename, mime):
        """读取上传文件二进制，直接存入数据库（跨部署永久保存，不依赖服务器硬盘）。"""
        data = file_obj.read()
        seq = self._row(self.conn.execute(
            "SELECT COALESCE(MAX(seq),0)+1 m FROM question_attachments WHERE question_id=?", (qid,)))["m"]
        cur = self.conn.execute(
            "INSERT INTO question_attachments (question_id, seq, filename, mime, stored_path, data, created_at) "
            "VALUES (?,?,?,?,?,?,?)",
            (qid, seq, filename, mime, "", data, _now()))
        self.conn.commit()
        return self._row(self.conn.execute(
            "SELECT att_id, question_id, seq, filename, mime, created_at "
            "FROM question_attachments WHERE att_id=?", (cur.lastrowid,)))

    def list_question_attachments(self, qid):
        return self._rows(self.conn.execute(
            "SELECT att_id, question_id, seq, filename, mime, created_at "
            "FROM question_attachments WHERE question_id=? ORDER BY seq", (qid,)))

    def get_attachment(self, att_id):
        return self._row(self.conn.execute("SELECT * FROM question_attachments WHERE att_id=?", (att_id,)))

    def delete_question_attachment(self, att_id):
        self.conn.execute("DELETE FROM question_attachments WHERE att_id=?", (att_id,))
        self.conn.commit()

    def _with_attachments(self, rows):
        """给题目列表/详情批量附加 attachments 字段。"""
        if not rows:
            return rows
        ids = [r["question_id"] for r in rows]
        placeholders = ",".join("?" * len(ids))
        atts = self._rows(self.conn.execute(
            f"SELECT att_id, question_id, seq, filename, mime, created_at "
            f"FROM question_attachments WHERE question_id IN ({placeholders}) ORDER BY seq", ids))
        by_q = {}
        for a in atts:
            by_q.setdefault(a["question_id"], []).append(a)
        for r in rows:
            r["attachments"] = by_q.get(r["question_id"], [])
        return rows

    # ---------- 题库 Excel 导入/导出 ----------
    def _norm_qtype(self, v):
        if not v:
            return "single"
        s = str(v).strip().lower()
        m = {"单选": "single", "single": "single", "多选": "multiple", "multiple": "multiple",
             "判断": "judge", "judge": "judge", "判断题": "judge", "简答": "essay",
             "essay": "essay", "问答": "essay"}
        return m.get(s, s)

    def _norm_bool(self, v):
        s = str(v).strip().lower()
        return "true" if s in ("true", "1", "正确", "对", "yes", "y") else "false"

    def import_questions_excel(self, path):
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        ws = wb.active
        header = [ws.cell(1, c).value for c in range(1, ws.max_column + 1)]

        def col(name):
            for i, h in enumerate(header):
                if h and name in str(h):
                    return i
            return None

        opt_cols = []
        for c in range(1, ws.max_column + 1):
            h = header[c - 1]
            if h and "选项" in str(h):
                opt_cols.append(c)
        cnt, errs = 0, []
        for r in range(2, ws.max_row + 1):
            ci = col("类型")
            typ = self._norm_qtype(ws.cell(r, ci + 1).value if ci is not None else None)
            cj = col("题面")
            content = ws.cell(r, cj + 1).value if cj is not None else None
            if not content or not str(content).strip():
                continue
            options = []
            if typ in ("single", "multiple", "judge"):
                keys = ["A", "B", "C", "D", "E", "F", "G"]
                for i, c in enumerate(opt_cols):
                    t = ws.cell(r, c).value
                    if t is not None and str(t).strip() != "":
                        options.append({"key": keys[i], "text": str(t).strip()})
                if typ == "judge":
                    options = [{"key": "true", "text": "正确"}, {"key": "false", "text": "错误"}]
            ca = col("正确答案")
            ans_raw = ws.cell(r, ca + 1).value if ca is not None else None
            answer = None
            if typ == "single":
                answer = str(ans_raw).strip().upper() if ans_raw is not None else None
            elif typ == "multiple":
                answer = json.dumps([x.strip().upper() for x in str(ans_raw).replace("，", ",").split(",") if x.strip()])
            elif typ == "judge":
                answer = self._norm_bool(ans_raw)
            cd = col("知识维度")
            dim_name = ws.cell(r, cd + 1).value if cd is not None else None
            dim_id = None
            if dim_name:
                d = self._row(self.conn.execute(
                    "SELECT dim_id FROM knowledge_dimensions WHERE name_cn=?", (str(dim_name).strip(),)))
                dim_id = d["dim_id"] if d else None
            cs = col("分值")
            score = ws.cell(r, cs + 1).value if cs is not None else None
            score = float(score) if score else 5
            ce = col("解析")
            explanation = ws.cell(r, ce + 1).value if ce is not None else None
            try:
                self.create_question({
                    "q_type": typ, "category": None, "content": str(content).strip(),
                    "options": json.dumps(options, ensure_ascii=False) if options else None,
                    "answer": answer, "dim_id": dim_id, "score": score, "explanation": explanation})
                cnt += 1
            except Exception as e:
                errs.append(f"第{r}行: {e}")
        return {"imported": cnt, "errors": errs}

    def export_questions_excel(self):
        import openpyxl
        from openpyxl.styles import Font
        qs = self.list_questions()
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "题库"
        head = ["类型", "题面", "选项A", "选项B", "选项C", "选项D", "选项E", "正确答案", "知识维度", "分值", "解析"]
        ws.append(head)
        for c in range(1, len(head) + 1):
            ws.cell(1, c).font = Font(bold=True)
        tmap = {"single": "单选", "multiple": "多选", "judge": "判断", "essay": "简答"}
        for q in qs:
            opts = json.loads(q["options"]) if q["options"] else []
            opt_map = {o["key"]: o["text"] for o in opts}
            if q["q_type"] == "multiple":
                ans = ",".join(json.loads(q["answer"])) if q["answer"] else ""
            else:
                ans = q["answer"] or ""
            ws.append([
                tmap.get(q["q_type"], q["q_type"]), q["content"],
                opt_map.get("A", ""), opt_map.get("B", ""), opt_map.get("C", ""),
                opt_map.get("D", ""), opt_map.get("E", ""),
                ans, q.get("dim_cn") or "", q["score"], q.get("explanation") or "",
            ])
        out = os.path.join(BASE_DIR, "data", "题库导出.xlsx")
        wb.save(out)
        return out

    # ---------- PPT 解析导入（考题原始 PPT：题干/选项/正确答案/解析） ----------
    def get_question_by_exam_q(self, exam_name, q_number):
        """按 (考试名, 题号) 取题库题目详情，用于弱项分析补充题干/答案/解析。"""
        row = self._row(self.conn.execute(
            "SELECT q.* FROM exam_question_bank b JOIN questions q ON b.question_id=q.question_id "
            "WHERE b.exam_name=? AND b.q_number=?", (exam_name, str(q_number))))
        if row:
            row["options_parsed"] = normalize_options(row["options"])
        return row

    def _parse_ppt_questions(self, path):
        """纯解析：把 PPT 拆成题目块列表（不落库）。
        关键过滤：只有「带 ≥2 个选项」或「已解析出答案」的块才视为题目，
        标题页/目录页/章节页/题干内嵌的编号步骤（无选项、无答案）会被丢弃，
        避免把非题目块误计为题目。返回 {questions:[...], skipped:int}。
        """
        from pptx import Presentation
        from pptx.util import Emu
        import re

        prs = Presentation(path)
        # 1) 收集全篇段落（幻灯片顺序；幻灯片内按阅读顺序 上→下、左→右）
        slides_paras = []
        for slide in prs.slides:
            shapes = sorted(
                [s for s in slide.shapes if s.has_text_frame],
                key=lambda s: (Emu(s.top).inches if s.top is not None else 0,
                               Emu(s.left).inches if s.left is not None else 0))
            paras = []
            for sh in shapes:
                for p in sh.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        paras.append(t)
            if paras:
                slides_paras.append(paras)

        # 2) 状态机解析
        questions = []
        skipped = 0
        cur = None
        in_feedback = False
        num_re = re.compile(r'^(\d{1,3})\s*[.、)]\s*(.*)$')
        opt_re = re.compile(r'^([A-Z])\s*[.、)]\s+(.+)$')
        type_markers = ["(t/f)", "(true/false)", "(multiple", "(single",
                        "(判断)", "(单选)", "(多选)"]

        def _looks_like_header(text):
            """极短的、或明显是封面/章节/目录/答案一览的编号块，不应算作题目。"""
            if not text:
                return False
            t = text.strip()
            if len(t) <= 6:
                return True
            # 答案一览单字母项：1) A / 1. A / 1、A / (1) A
            if re.match(r'^\d+\s*[.、)]\s*[A-Z]\s*$', t):
                return True
            if re.match(r'^[（(]\d+[)）]\s*[A-Z]\s*$', t):
                return True
            hdr = re.compile(r'^(第[一二三四五六七八九十\d]+\s*章|chapter|agenda|目录|contents|'
                             r'overview|封面|summary|quiz\s*$|section\s*\d|'
                             r'answers?$|answer\s*key|correct\s*answers?|key\s*$|'
                             r'答案$|答案\s*key|参考答案|答题|解析)', re.I)
            return bool(hdr.match(t))

        # 题型标记（不区分大小写、忽略空格），用于区分「真题目」与「编号步骤/目录」
        _qtype_hints = ["multiple choice", "single choice", "t/f", "true/false",
                        "choose all", "y/n", "free text", "单选", "判断", "多选",
                        "(t/f)", "（判断）", "（单选）", "（多选）", "true or false"]
        def _has_qmarker(text):
            tl = (text or "").lower()
            return any(h in tl for h in _qtype_hints)

        # 非字母选项识别：中文 正确/错误/对/错/是/否、符号 √/×、编号子选项 ①/⑵/(1)/1)
        _opt_zh_re = re.compile(r'^(正确|错误|对|错|是|否|√|×|✓|✗)\b', re.I)
        _opt_num_re = re.compile(r'^[①-⑨⑩⑪-⑳]|^[（(]\d+[)）]|^\d+[)）、]')
        def _is_option_like(text):
            if opt_re.match(text):
                return True
            s = text.strip()
            if _opt_zh_re.match(s):
                return True
            if _opt_num_re.match(s):
                return True
            return False

        def finish(c):
            nonlocal skipped
            if not c or not c.get("content"):
                return
            # 保留条件（收紧，修复 Order PPT 多解析 5 道的问题）：
            # ① 选项≥2；② 已解析出答案；③ 是「编号题干起点」且本身像题目：
            #    须带题型标记（Multiple Choice/T/F/单选…）、问号，或至少 1 个选项标记
            #    （含中文 正确/错误、符号 √/×、编号子选项）。
            # 目录页/章节页/答案一览/编号步骤（无选项、无答案、无题型标记）一律丢弃。
            has_opts = len(c.get("options") or []) >= 2
            has_ans = bool(c.get("answer"))
            opt_tokens = c.get("_opt_tokens", 0)
            is_numbered_q = bool(c.get("_numbered")) and not _looks_like_header(c.get("content"))
            qmarker = _has_qmarker(c.get("content"))
            qmark = "?" in (c.get("content") or "")
            if has_opts or has_ans or (is_numbered_q and (qmarker or qmark or opt_tokens >= 1)):
                questions.append(c)
            else:
                skipped += 1

        for paras in slides_paras:
            for t in paras:
                tl = t.lower()
                mnum = num_re.match(t)
                is_start = bool(mnum)
                if not is_start:
                    is_start = (any(mk in tl for mk in type_markers)
                                and not tl.startswith("answer")
                                and not tl.startswith("feedback"))
                if is_start:
                    finish(cur)
                    cur = {"content": t, "options": [], "answer": None, "explanation": "",
                           "_await_answer": False, "_feedback_open": False,
                           "_numbered": bool(mnum), "_opt_tokens": 0}
                    in_feedback = False
                    continue
                if cur is None:
                    continue
                # 选项（仅在该题答案尚未出现前收集）
                mopt = opt_re.match(t)
                if mopt and not in_feedback and cur["answer"] is None:
                    cur["options"].append({"key": mopt.group(1), "text": mopt.group(2).strip()})
                    cur["_opt_tokens"] = cur.get("_opt_tokens", 0) + 1
                    continue
                # 非字母选项（中文 正确/错误、符号 √/×、编号子选项 ①/⑵/(1) 等）：
                # 计入选项，使「选项≥2」即可保留；同时避免把编号步骤误当题目。
                if (not in_feedback) and cur["answer"] is None and _is_option_like(t):
                    cur["_opt_tokens"] = cur.get("_opt_tokens", 0) + 1
                    cur["options"].append({"key": f"o{cur['_opt_tokens']}", "text": t.strip()})
                    continue
                # 正确答案（兼容英文 answer: 与中文 答案：/正确答案：）
                if ("answer" in tl or "答案" in t or "正确答案" in t) and not in_feedback:
                    val = re.sub(r'^.*(?:answer|答案|正确答案)\s*[:：]\s*', '', t, flags=re.I).strip().rstrip('. ')
                    if val:
                        cur["answer"] = val
                        in_feedback = False
                        continue
                    else:
                        cur["_await_answer"] = True
                        continue
                if cur.get("_await_answer"):
                    cand = t.strip().rstrip('. ')
                    if cand:
                        cur["answer"] = cand
                        cur["_await_answer"] = False
                        continue
                # 解析
                if "feedback" in tl:
                    cur["explanation"] = re.sub(r'^.*feedback\s*[:：]\s*', '', t, flags=re.I).strip()
                    cur["_feedback_open"] = True
                    in_feedback = True
                    continue
                if in_feedback or cur.get("_feedback_open"):
                    cur["explanation"] = (cur["explanation"] + " " + t).strip() if cur["explanation"] else t
                    continue
                # 其它续行（选项换行等）：并入最后一条选项
                if cur["options"]:
                    cur["options"][-1]["text"] += " " + t
        finish(cur)
        return {"questions": questions, "skipped": skipped}

    def preview_questions_ppt(self, path):
        """预览解析结果（不落库），供上传前确认题数与内容。"""
        p = self._parse_ppt_questions(path)
        items = []
        for i, q in enumerate(p["questions"]):
            items.append({
                "seq": i + 1,
                "content": (q["content"] or "")[:140],
                "options": len(q.get("options") or []),
                "answer": q.get("answer"),
                "q_type": self._detect_ppt_qtype(q["content"], q.get("answer") or ""),
            })
        return {"count": len(p["questions"]), "skipped": p["skipped"], "questions": items}

    def import_questions_ppt(self, path, exam_name, dim_ids=None, orig_filename=None):
        """解析考题 PPT（题目+选项+正确答案+解析），写入题库并关联到指定考试。
        支持版式：题干 '1. (Multiple Choice) ...' / '7. (T/F) ...'；
        选项 'A. ...'；正确答案显式 'Answer: A'（绿色标记）；解析 'Feedback: ...'。
        题号按解析顺序 Q1..QN 与 exam_results.subjects 的键对齐。
        dim_ids: 维度 ID 列表，支持一道题同时归属多个维度（自由组合）。
        返回 {exam_name, count, questions:[...], warning}
        """
        import re
        parsed = self._parse_ppt_questions(path)
        questions = parsed["questions"]
        skipped = parsed["skipped"]

        # 3) 规范化并落库
        exam_name = re.sub(r"\s+", " ", exam_name).strip()
        # 重新导入前清掉该考试旧的题库/题目/维度映射，保证题号对齐、无残留旧题
        self.conn.execute("DELETE FROM exam_question_dimensions WHERE exam_name=?", (exam_name,))
        self.conn.execute("DELETE FROM exam_question_bank WHERE exam_name=?", (exam_name,))
        self.conn.execute(
            "DELETE FROM questions WHERE category=? "
            "AND question_id NOT IN (SELECT question_id FROM paper_questions)",
            (exam_name,))
        qmax = self._exam_qmax(exam_name)
        res_qcount = self._exam_result_qcount(exam_name)
        out = []
        warn = None
        dim_list = [int(x) for x in (dim_ids or [])]
        for i, q in enumerate(questions):
            q_number = f"Q{i + 1}"
            ans = (q["answer"] or "").strip()
            qtype = self._detect_ppt_qtype(q["content"], ans)
            if qtype == "judge":
                options = [{"key": "true", "text": "正确"}, {"key": "false", "text": "错误"}]
                norm_ans = "true" if ans.lower().startswith("t") else ("false" if ans.lower().startswith("f") else None)
            elif qtype == "multiple":
                letters = re.findall(r'[A-Z]', ans)
                options = q["options"] or []
                norm_ans = json.dumps([x.upper() for x in letters])
            else:
                options = q["options"] or []
                norm_ans = ans[:1].upper() if ans else None
            opts_json = json.dumps(options, ensure_ascii=False) if options else None
            # questions.dim_id 仅作单值回退（题目库筛选用）；真正的多维度以 exam_question_dimensions 为准
            primary_dim = dim_list[0] if dim_list else None
            existing = self._row(self.conn.execute(
                "SELECT question_id FROM exam_question_bank WHERE exam_name=? AND q_number=?",
                (exam_name, q_number)))
            if existing:
                self.conn.execute(
                    "UPDATE questions SET q_type=?, content=?, options=?, answer=?, dim_id=?, "
                    "score=?, explanation=? WHERE question_id=?",
                    (qtype, q["content"].strip(), opts_json, norm_ans, primary_dim, 5.0,
                     q["explanation"].strip() or None, existing["question_id"]))
                question_id = existing["question_id"]
            else:
                cur = self.conn.execute(
                    "INSERT INTO questions (q_type, category, content, options, answer, dim_id, score, explanation, created_at) "
                    "VALUES (?,?,?,?,?,?,?,?,?)",
                    (qtype, exam_name, q["content"].strip(), opts_json, norm_ans,
                     primary_dim, 5.0, q["explanation"].strip() or None, _now()))
                question_id = cur.lastrowid
            self.conn.execute(
                "INSERT OR REPLACE INTO exam_question_bank (exam_name, q_number, question_id, seq) "
                "VALUES (?,?,?,?)", (exam_name, q_number, question_id, i + 1))
            # 一题可挂多个维度（自由组合）：先清旧映射，再按勾选写入
            self.conn.execute(
                "DELETE FROM exam_question_dimensions WHERE exam_name=? AND q_number=?",
                (exam_name, q_number))
            for did in dim_list:
                self.set_question_dimension(exam_name, q_number, int(did),
                                            max_score=qmax.get(q_number))
            out.append({"q_number": q_number, "content": q["content"], "q_type": qtype,
                        "answer": norm_ans, "options": options, "dim_ids": dim_list})
        # 记录/更新该考试的 PPT 元数据，供「考题 PPT 管理」列出与删除
        self.conn.execute(
            "INSERT OR REPLACE INTO question_bank_meta "
            "(exam_name, orig_filename, uploaded_at, q_count, dim_ids) VALUES (?,?,?,?,?)",
            (exam_name, orig_filename or None, _now(), len(questions),
             ",".join(str(d) for d in dim_list) if dim_list else None))
        self.conn.commit()
        if skipped:
            note = f"（已自动忽略 {skipped} 个非题目块：目录/章节/编号步骤等）"
        else:
            note = ""
        if res_qcount and res_qcount != len(questions):
            warn = (f"本次解析并按题序导入 {len(questions)} 道题{note}；"
                    f"但该考试成绩里共有 {res_qcount} 道小题——两者不一致。\n"
                    f"请核对：是否 PPT 比考试多/少一题，或题序错位。弱项分析按题号 Q1..QN 一一对应，"
                    f"只要顺序一致，前 {min(len(questions), res_qcount)} 题可正确关联，多出的题无害。")
        elif skipped:
            warn = f"已导入 {len(questions)} 道题{note}。"
        return {"exam_name": exam_name, "count": len(questions), "skipped": skipped,
                "questions": out, "warning": warn}

    def _detect_ppt_qtype(self, content, answer):
        cl = (content or "").lower()
        if "(t/f)" in cl or "(true/false)" in cl or "(判断)" in cl:
            return "judge"
        if answer and len(re.findall(r'[A-Z]', answer)) > 1:
            return "multiple"
        return "single"

    def _exam_qmax(self, exam_name):
        rows = self._rows(self.conn.execute(
            "SELECT r.subjects FROM exam_results r JOIN exam_sessions s ON r.session_id=s.session_id "
            "WHERE s.exam_name=?", (exam_name,)))
        qmax = {}
        for r in rows:
            subj = json.loads(r["subjects"]) if r["subjects"] else {}
            for k, v in subj.items():
                try:
                    v = float(v)
                except (TypeError, ValueError):
                    continue
                qmax[k] = max(qmax.get(k, 0), v)
        return qmax

    def _exam_result_qcount(self, exam_name):
        # 优先用会话记录的小题数（表头识别，含自由作答等无分题）；
        # 旧导入或缺失时退化为按成绩 subject 键数取最大值。
        noted = 0
        for s in self._rows(self.conn.execute(
                "SELECT note FROM exam_sessions WHERE exam_name=?", (exam_name,))):
            m = re.search(r"共\s*(\d+)\s*道小题", s.get("note") or "")
            if m:
                noted = max(noted, int(m.group(1)))
        mx = 0
        for r in self._rows(self.conn.execute(
                "SELECT r.subjects FROM exam_results r JOIN exam_sessions s "
                "ON r.session_id=s.session_id WHERE s.exam_name=?", (exam_name,))):
            subj = json.loads(r["subjects"]) if r["subjects"] else {}
            mx = max(mx, len(subj))
        return max(noted, mx)

    # ---------- 考题 PPT / 题库管理 ----------
    def list_question_banks(self):
        """列出所有已上传考题 PPT（按考试聚合）。
        对每个考试返回：文件名、上传时间、题库题数、成绩小题数、是否已关联成绩、
        覆盖的维度名列表、报名/答卷份数。用于「考题 PPT 管理」面板的列出与删除。
        兼容早期未记录元数据的历史导入（从 exam_question_bank 回填）。"""
        # 所有出现过的考试名（题库表 ∪ 元数据表）
        names = set()
        for r in self._rows(self.conn.execute(
                "SELECT DISTINCT exam_name FROM exam_question_bank")):
            names.add(r["exam_name"])
        for r in self._rows(self.conn.execute(
                "SELECT exam_name FROM question_bank_meta")):
            names.add(r["exam_name"])
        out = []
        for name in sorted(names):
            meta = self._row(self.conn.execute(
                "SELECT * FROM question_bank_meta WHERE exam_name=?", (name,)))
            q_count = self._row(self.conn.execute(
                "SELECT COUNT(*) c FROM exam_question_bank WHERE exam_name=?", (name,)))["c"]
            # 覆盖到的维度名（去重）
            dims = self._rows(self.conn.execute(
                "SELECT DISTINCT d.dim_id, d.name_cn, d.name_en "
                "FROM exam_question_dimensions e JOIN knowledge_dimensions d ON e.dim_id=d.dim_id "
                "WHERE e.exam_name=? ORDER BY d.dim_id", (name,)))
            # 是否已关联成绩（存在同名考试批次）
            sess = self._row(self.conn.execute(
                "SELECT COUNT(*) c FROM exam_sessions WHERE exam_name=?", (name,)))["c"]
            result_qcount = self._exam_result_qcount(name)
            out.append({
                "exam_name": name,
                "orig_filename": (meta or {}).get("orig_filename"),
                "uploaded_at": (meta or {}).get("uploaded_at"),
                "q_count": q_count,
                "result_qcount": result_qcount,
                "linked": bool(sess),
                "dims": [{"dim_id": d["dim_id"], "name_cn": d["name_cn"],
                          "name_en": d["name_en"]} for d in dims],
            })
        return out

    def delete_question_bank(self, exam_name):
        """删除某考试对应的整套上传题库（PPT 关联）：
        清 question_bank_meta / exam_question_dimensions / exam_question_bank，
        并删除仅属于该题库、未被任何试卷引用的题目。不影响考试成绩本身。"""
        exam_name = (exam_name or "").strip()
        if not exam_name:
            raise ValueError("exam_name 不能为空")
        # 先记下本题库引用的题目 ID（因 exam_question_bank 外键指向 questions，需先删关联表）
        qids = [r["question_id"] for r in self._rows(self.conn.execute(
            "SELECT question_id FROM exam_question_bank WHERE exam_name=?", (exam_name,)))]
        self.conn.execute("DELETE FROM exam_question_dimensions WHERE exam_name=?", (exam_name,))
        self.conn.execute("DELETE FROM exam_question_bank WHERE exam_name=?", (exam_name,))
        self.conn.execute("DELETE FROM question_bank_meta WHERE exam_name=?", (exam_name,))
        # 再删除仅属于该题库、且未被任何试卷引用的题目（避免误删组卷题）
        for qid in qids:
            used = self._row(self.conn.execute(
                "SELECT 1 FROM paper_questions WHERE question_id=? LIMIT 1", (qid,)))
            if not used:
                self.conn.execute(
                    "DELETE FROM questions WHERE question_id=? AND category=?",
                    (qid, exam_name))
        self.conn.commit()
        return {"exam_name": exam_name}

    def bulk_delete_question_banks(self, names: list) -> int:
        """批量删除整套上传题库（按 exam_name）。返回成功删除的套数。"""
        n = 0
        for name in (names or []):
            try:
                self.delete_question_bank(name)
                n += 1
            except Exception:
                pass
        return n

    # ---------- 试卷 / 组卷 ----------
    def list_papers(self, filters=None):
        filters = filters or {}
        sql = "SELECT * FROM exam_papers WHERE 1=1"
        vals = []
        if filters.get("status"):
            sql += " AND status=?"; vals.append(filters["status"])
        sql += " ORDER BY paper_id DESC"
        rows = self._rows(self.conn.execute(sql, vals))
        for r in rows:
            r["question_count"] = self._row(self.conn.execute(
                "SELECT COUNT(*) c FROM paper_questions WHERE paper_id=?", (r["paper_id"],)))["c"]
        return rows

    def get_paper(self, paper_id):
        p = self._row(self.conn.execute("SELECT * FROM exam_papers WHERE paper_id=?", (paper_id,)))
        if p:
            p["assignments"] = self.get_paper_assignments(paper_id)
        return p

    def create_paper(self, data):
        cur = self.conn.execute(
            "INSERT INTO exam_papers (title, batch, exam_type, status, duration_min, pass_score, created_by, created_at) "
            "VALUES (?,?,?,?,?,?,?,?)",
            (data["title"].strip(), data.get("batch", "online"), data.get("exam_type", "onboarding"),
             "draft", int(data.get("duration_min") or 0), float(data.get("pass_score") or 60),
             data.get("created_by"), _now()))
        self.conn.commit()
        return self.get_paper(cur.lastrowid)

    def update_paper(self, paper_id, data):
        fields, vals = [], []
        for k in ("title", "batch", "exam_type", "duration_min", "pass_score", "status", "open_at", "close_at"):
            if k in data:
                v = data[k]
                if k == "duration_min":
                    v = int(v or 0)
                elif k == "pass_score":
                    v = float(v or 60)
                elif isinstance(v, str):
                    v = v.strip()
                fields.append(f"{k}=?"); vals.append(v)
        if fields:
            vals.append(paper_id)
            self.conn.execute(f"UPDATE exam_papers SET {','.join(fields)} WHERE paper_id=?", vals)
            self.conn.commit()
        return self.get_paper(paper_id) or {}

    def delete_paper(self, paper_id):
        self.conn.execute(
            "DELETE FROM exam_answers WHERE attempt_id IN (SELECT attempt_id FROM exam_attempts WHERE paper_id=?)",
            (paper_id,))
        self.conn.execute("DELETE FROM exam_attempts WHERE paper_id=?", (paper_id,))
        self.conn.execute("DELETE FROM paper_questions WHERE paper_id=?", (paper_id,))
        self.conn.execute("DELETE FROM exam_assignments WHERE paper_id=?", (paper_id,))
        self.conn.execute("DELETE FROM exam_papers WHERE paper_id=?", (paper_id,))
        self.conn.commit()

    def publish_paper(self, paper_id, open_at=None, close_at=None):
        now = _now()
        self.conn.execute(
            "UPDATE exam_papers SET status='published', open_at=?, close_at=? WHERE paper_id=?",
            (open_at or now, close_at, paper_id))
        self.conn.commit()
        return self.get_paper(paper_id) or {}

    def get_paper_questions(self, paper_id):
        rows = self._rows(self.conn.execute(
            "SELECT pq.seq, pq.score, q.*, d.name_cn AS dim_cn, d.name_en AS dim_en "
            "FROM paper_questions pq JOIN questions q ON pq.question_id=q.question_id "
            "LEFT JOIN knowledge_dimensions d ON q.dim_id=d.dim_id "
            "WHERE pq.paper_id=? ORDER BY pq.seq", (paper_id,)))
        for r in rows:
            r["options_parsed"] = normalize_options(r["options"])
        return self._with_attachments(rows)

    def set_paper_questions(self, paper_id, items):
        self.conn.execute("DELETE FROM paper_questions WHERE paper_id=?", (paper_id,))
        for i, it in enumerate(items):
            self.conn.execute(
                "INSERT INTO paper_questions (paper_id, question_id, seq, score) VALUES (?,?,?,?)",
                (paper_id, int(it["question_id"]), it.get("seq", i + 1), self._paper_question_score(it)))
        self.conn.commit()

    def _paper_question_score(self, it):
        """组卷每题分值：优先用前端传入分；否则取题库该题真实分值；兜底 5。"""
        s = it.get("score")
        if s not in (None, "", 0):
            try:
                return float(s)
            except (TypeError, ValueError):
                pass
        q = self._row(self.conn.execute(
            "SELECT score FROM questions WHERE question_id=?", (int(it["question_id"]),)))
        return float(q["score"]) if (q and q.get("score")) else 5.0

    # ---------- 考试分配（指定客服可见） ----------
    def set_paper_assignments(self, paper_id, rep_ids):
        """设置试卷可见的客服列表。rep_ids 为空列表 = 全员广播（不限人）。"""
        self.conn.execute("DELETE FROM exam_assignments WHERE paper_id=?", (paper_id,))
        for rid in (rep_ids or []):
            rid = (rid or "").strip()
            if not rid:
                continue
            self.conn.execute("INSERT OR IGNORE INTO exam_assignments (paper_id, rep_id, created_at) VALUES (?,?,?)",
                              (paper_id, rid, _now()))
        self.conn.commit()

    def get_paper_assignments(self, paper_id):
        return [r["rep_id"] for r in self._rows(self.conn.execute(
            "SELECT rep_id FROM exam_assignments WHERE paper_id=?", (paper_id,)))]

    # ---------- 考试 / 判分 ----------
    def _rep_in_assignment(self, paper_id, rep_id):
        """试卷无任何分配=全员可见；否则仅当 rep 在分配名单中。"""
        rows = self._rows(self.conn.execute(
            "SELECT rep_id FROM exam_assignments WHERE paper_id=?", (paper_id,)))
        if not rows:
            return True
        return any(r["rep_id"] == rep_id for r in rows)

    def list_available_papers(self, rep_id):
        now = _now()
        papers = self._rows(self.conn.execute(
            "SELECT p.* FROM exam_papers p WHERE p.status='published' ORDER BY p.paper_id DESC", ()))
        out = []
        for p in papers:
            paper_ok = (p["open_at"] is None or p["open_at"] <= now) and (p["close_at"] is None or p["close_at"] >= now)
            global_visible = paper_ok and self._rep_in_assignment(p["paper_id"], rep_id)
            # 每客服独立补考窗口
            row = self._row(self.conn.execute(
                "SELECT open_at, due_at FROM exam_assignments WHERE paper_id=? AND rep_id=?", (p["paper_id"], rep_id)))
            # 仅当该客服确实有独立时间窗口时才按个人窗口放行；
            # 普通分配（窗口为 NULL）或无窗口（已撤销补考）一律回退到试卷全局规则。
            personal_ok = bool(row) and (row["open_at"] is not None or row["due_at"] is not None) \
                and (row["open_at"] is None or row["open_at"] <= now) and (row["due_at"] is None or row["due_at"] >= now)
            if not (global_visible or personal_ok):
                continue
            taken = self._row(self.conn.execute(
                "SELECT attempt_id, status FROM exam_attempts WHERE paper_id=? AND rep_id=? "
                "ORDER BY attempt_id DESC LIMIT 1", (p["paper_id"], rep_id)))
            p["already_taken"] = bool(taken)
            p["attempt_status"] = taken["status"] if taken else None
            out.append(p)
        return out

    def start_attempt(self, paper_id, rep_id):
        paper = self.get_paper(paper_id)
        if not paper or paper["status"] != "published":
            raise ValueError("试卷未发布或不存在")
        now = _now()
        row = self._row(self.conn.execute(
            "SELECT open_at, due_at FROM exam_assignments WHERE paper_id=? AND rep_id=?", (paper_id, rep_id)))
        # 仅当该客服确有独立时间窗口时才按个人窗口放行；否则回退到试卷全局规则。
        personal_ok = bool(row) and (row["open_at"] is not None or row["due_at"] is not None) \
            and (row["open_at"] is None or row["open_at"] <= now) and (row["due_at"] is None or row["due_at"] >= now)
        global_ok = (paper["open_at"] is None or paper["open_at"] <= now) and (paper["close_at"] is None or paper["close_at"] >= now)
        if not (global_ok or personal_ok):
            raise ValueError("考试已截止（含补考窗口）")
        ex = self._row(self.conn.execute(
            "SELECT * FROM exam_attempts WHERE paper_id=? AND rep_id=? AND status='in_progress' "
            "ORDER BY attempt_id DESC LIMIT 1", (paper_id, rep_id)))
        if ex:
            return ex
        cur = self.conn.execute(
            "INSERT INTO exam_attempts (paper_id, rep_id, start_time, status, created_at) "
            "VALUES (?,?,?,?,?)",
            (paper_id, rep_id, now, "in_progress", now))
        self.conn.commit()
        return self._row(self.conn.execute("SELECT * FROM exam_attempts WHERE attempt_id=?", (cur.lastrowid,)))

    def _norm_keys(self, val):
        """把答案（'A' / 'A,C' / '["A","C"]' / ['A','C']）统一成大写字母集合，便于多选题比较。"""
        if val is None:
            return set()
        if isinstance(val, str):
            s = val.strip()
            if s.startswith("["):
                try:
                    return set(str(x).strip().upper() for x in json.loads(s))
                except Exception:
                    pass
            return set(x.strip().upper() for x in s.replace("，", ",").split(",") if x.strip())
        if isinstance(val, (list, tuple, set)):
            return set(str(x).strip().upper() for x in val)
        return set([str(val).strip().upper()])

    def _grade_one(self, question, student_answer):
        qtype = question["q_type"]
        if qtype == "essay":
            return None, 0.0
        if qtype == "single":
            correct = str(student_answer).strip().upper() in self._norm_keys(question["answer"])
        elif qtype == "multiple":
            correct = (self._norm_keys(student_answer) == self._norm_keys(question["answer"]))
        elif qtype == "judge":
            correct = self._norm_bool(student_answer) == str(question["answer"]).strip().lower()
        else:
            correct = False
        return (1 if correct else 0), (float(question["score"]) if correct else 0.0)

    def submit_attempt(self, attempt_id, answers):
        attempt = self._row(self.conn.execute("SELECT * FROM exam_attempts WHERE attempt_id=?", (attempt_id,)))
        if not attempt:
            raise ValueError("考试会话不存在")
        paper_id = attempt["paper_id"]
        pq = self.get_paper_questions(paper_id)
        auto_total = 0.0
        self.conn.execute("DELETE FROM exam_answers WHERE attempt_id=?", (attempt_id,))
        for q in pq:
            qid = q["question_id"]
            sa = answers.get(str(qid))
            is_correct, sc = self._grade_one(q, sa)
            self.conn.execute(
                "INSERT INTO exam_answers (attempt_id, question_id, answer, is_correct, score) VALUES (?,?,?,?,?)",
                (attempt_id, qid, json.dumps(sa) if not isinstance(sa, str) else sa, is_correct, sc))
            if is_correct == 1:
                auto_total += sc
        now = _now()
        has_essay = any(q["q_type"] == "essay" for q in pq)
        status = "submitted" if has_essay else "graded"
        pscore = self.get_paper(paper_id)["pass_score"] or 60
        ratio = self.get_pass_line_ratio()
        full_score = self._paper_full_score(paper_id)
        # 在线交卷本就是在线路径：必须用百分比判定（online=True），
        # 否则会落到非在线分支按「绝对分 total>=pass_score」判定，
        # 满分低的在线考试（如满分10得8分）会被误判未通过。
        score_rate, passed = self._score_rate_pass(round(auto_total, 2), full_score, pscore, ratio, online=True)
        self.conn.execute(
            "UPDATE exam_attempts SET submit_time=?, status=?, auto_score=?, total_score=?, passed=?, score_rate=?, full_score=? WHERE attempt_id=?",
            (now, status, round(auto_total, 2), round(auto_total, 2), passed, score_rate, full_score, attempt_id))
        self.conn.commit()
        self._sync_attempt_to_results(attempt_id)
        return self.get_attempt(attempt_id)

    def _sync_attempt_to_results(self, attempt_id):
        attempt = self._row(self.conn.execute("SELECT * FROM exam_attempts WHERE attempt_id=?", (attempt_id,)))
        paper = self.get_paper(attempt["paper_id"])
        rep = self._row(self.conn.execute("SELECT name FROM cs_reps WHERE rep_id=?", (attempt["rep_id"],)))
        rep_name = rep["name"] if rep else attempt["rep_id"]
        pq = self.get_paper_questions(attempt["paper_id"])
        ans = self._rows(self.conn.execute("SELECT * FROM exam_answers WHERE attempt_id=?", (attempt_id,)))
        ans_map = {a["question_id"]: a for a in ans}
        session = self._row(self.conn.execute(
            "SELECT * FROM exam_sessions WHERE exam_name=? AND note=?",
            (paper["title"], f"online:paper_id={paper['paper_id']}")))
        if not session:
            session = self.create_session({
                "exam_name": paper["title"], "batch": paper.get("batch") or "online",
                "exam_date": attempt["submit_time"][:10] if attempt["submit_time"] else _now()[:10],
                "pass_score": paper["pass_score"], "exam_type": paper["exam_type"],
                "note": f"online:paper_id={paper['paper_id']}"})
        subjects = {}
        for q in pq:
            a = ans_map.get(q["question_id"])
            seq = q["seq"] or 1
            subjects[f"Q{seq}"] = a["score"] if a else 0
            # 同步写入 exam_question_bank，使「弱项分析」能取到题干/答案/解析
            # （无论试卷是 PPT 导入还是题库手动组卷，弱项展示都能显示真实错题）。
            # 用 INSERT OR IGNORE：若该题已来自 PPT 导入则保留原映射，避免覆盖。
            self.conn.execute(
                "INSERT OR IGNORE INTO exam_question_bank (exam_name, q_number, question_id, seq) VALUES (?,?,?,?)",
                (paper["title"], f"Q{seq}", q["question_id"], seq))
            if q["dim_id"]:
                self.set_question_dimension(paper["title"], f"Q{seq}", q["dim_id"],
                                            max_score=q.get("score") or q["score"])
        exist = self._row(self.conn.execute(
            "SELECT result_id FROM exam_results WHERE session_id=? AND rep_id=?",
            (session["session_id"], attempt["rep_id"])))
        total = attempt["total_score"]
        passed = attempt["passed"]
        score_rate = attempt.get("score_rate")
        full_score = attempt.get("full_score")
        if exist:
            self.conn.execute(
                "UPDATE exam_results SET name=?, subjects=?, total=?, passed=?, score_rate=?, full_score=? WHERE result_id=?",
                (rep_name, json.dumps(subjects, ensure_ascii=False), total, passed, score_rate, full_score, exist["result_id"]))
        else:
            self.conn.execute(
                "INSERT INTO exam_results (session_id, rep_id, name, subjects, total, passed, score_rate, full_score, created_at) "
                "VALUES (?,?,?,?,?,?,?,?,?)",
                (session["session_id"], attempt["rep_id"], rep_name,
                 json.dumps(subjects, ensure_ascii=False), total, passed, score_rate, full_score, _now()))
        self.conn.commit()
        # 积分：不通过无分；通过只取分层最高档（不叠参与基础分）
        try:
            rid = attempt["rep_id"]
            if passed == 1:
                pd = self.pass_points_for(score_rate, passed=True)
                if pd:
                    self.award_points(rid, "pass", "attempt", attempt_id, "通过在线考试", delta=pd)
        except Exception as _e:
            print(f"[points] sync award skipped: {_e}")

    def get_attempt(self, attempt_id):
        attempt = self._row(self.conn.execute("SELECT * FROM exam_attempts WHERE attempt_id=?", (attempt_id,)))
        if not attempt:
            return {}
        paper = self.get_paper(attempt["paper_id"])
        pq = self.get_paper_questions(attempt["paper_id"])
        ans = self._rows(self.conn.execute("SELECT * FROM exam_answers WHERE attempt_id=?", (attempt_id,)))
        ans_map = {a["question_id"]: a for a in ans}
        detail = []
        for q in pq:
            a = ans_map.get(q["question_id"])
            detail.append({
                "question_id": q["question_id"], "seq": q["seq"], "q_type": q["q_type"],
                "content": q["content"], "options": q["options_parsed"],
                "answer": q["answer"], "dim_cn": q.get("dim_cn"),
                "student_answer": a["answer"] if a else None,
                "is_correct": a["is_correct"] if a else None,
                "score": a["score"] if a else 0,
                "graded_by": a["graded_by"] if a else None,
                "explanation": q["explanation"],
            })
        return {"attempt": attempt, "paper": paper, "detail": detail}

    def list_attempts(self, filters=None):
        filters = filters or {}
        sql = ("SELECT a.*, p.title, r.name FROM exam_attempts a "
               "JOIN exam_papers p ON a.paper_id=p.paper_id "
               "JOIN cs_reps r ON a.rep_id=r.rep_id WHERE 1=1")
        vals = []
        if filters.get("paper_id"):
            sql += " AND a.paper_id=?"; vals.append(int(filters["paper_id"]))
        if filters.get("rep_id"):
            sql += " AND a.rep_id=?"; vals.append(filters["rep_id"])
        if filters.get("status"):
            sql += " AND a.status=?"; vals.append(filters["status"])
        sql += " ORDER BY a.attempt_id DESC"
        return self._rows(self.conn.execute(sql, vals))

    def my_exam_history(self, rep_id: str) -> list:
        """客服端「我的考试历史」：返回该客服全部考试会话（exam_attempts），
        无论对应试卷当前是否仍 published（已关闭/下线的历史考试也列出）。
        每条附带其对应的 exam_sessions.session_id（用于维度弱项按场分析）。"""
        attempts = self._rows(self.conn.execute(
            "SELECT a.*, p.title, p.pass_score FROM exam_attempts a "
            "JOIN exam_papers p ON a.paper_id=p.paper_id "
            "WHERE a.rep_id=? ORDER BY a.attempt_id DESC", (rep_id,)))
        out = []
        for a in attempts:
            session = self._row(self.conn.execute(
                "SELECT session_id FROM exam_sessions WHERE note=?",
                (f"online:paper_id={a['paper_id']}",)))
            out.append({
                "attempt_id": a["attempt_id"],
                "paper_id": a["paper_id"],
                "paper_title": a["title"],
                "pass_score": a["pass_score"],
                "start_time": a["start_time"],
                "submit_time": a["submit_time"],
                "total_score": a["total_score"],
                "passed": a["passed"],
                "status": a["status"],
                "session_id": session["session_id"] if session else None,
            })
        return out

    def list_pending_grading(self):
        return self._rows(self.conn.execute(
            "SELECT ea.attempt_id, ea.question_id, ea.answer, a.rep_id, r.name AS rep_name, "
            "p.title AS paper_title, q.content, q.score AS max_score, a.paper_id "
            "FROM exam_answers ea "
            "JOIN exam_attempts a ON ea.attempt_id=a.attempt_id "
            "JOIN exam_papers p ON a.paper_id=p.paper_id "
            "JOIN questions q ON ea.question_id=q.question_id "
            "JOIN cs_reps r ON a.rep_id=r.rep_id "
            "WHERE q.q_type='essay' AND ea.is_correct IS NULL "
            "ORDER BY ea.attempt_id DESC"))

    def grade_essay(self, attempt_id, question_id, score, grader):
        attempt = self._row(self.conn.execute("SELECT * FROM exam_attempts WHERE attempt_id=?", (attempt_id,)))
        q = self._row(self.conn.execute("SELECT * FROM questions WHERE question_id=?", (question_id,)))
        max_score = q["score"] if q else 0
        sc = max(0.0, min(float(score), float(max_score)))
        now = _now()
        self.conn.execute(
            "UPDATE exam_answers SET is_correct=?, score=?, graded_by=?, graded_at=? "
            "WHERE attempt_id=? AND question_id=?",
            (1 if sc >= max_score else 0, sc, grader, now, attempt_id, question_id))
        manual = self._row(self.conn.execute(
            "SELECT SUM(score) s FROM exam_answers WHERE attempt_id=? AND question_id IN "
            "(SELECT question_id FROM questions WHERE q_type='essay')", (attempt_id,)))["s"] or 0
        auto = self._row(self.conn.execute(
            "SELECT auto_score FROM exam_attempts WHERE attempt_id=?", (attempt_id,)))["auto_score"] or 0
        total = auto + manual
        paper = self.get_paper(attempt["paper_id"])
        full_score = self._paper_full_score(attempt["paper_id"])
        ratio = self.get_pass_line_ratio()
        # 在线考试按该批次(paper 对应 session)的 pass_score 作百分比阈值判定，
        # 否则会走非在线分支把 total 当绝对分与 60 比（如 9>=60 误判未通过）；
        # 与 recompute_score_rates 的判定保持一致，确保判分前后/个人视图与管理端一致。
        sess = self._row(self.conn.execute(
            "SELECT * FROM exam_sessions WHERE note=?", (f"online:paper_id={attempt['paper_id']}",)))
        if sess:
            ps = sess["pass_score"]
            online = True
        else:
            ps = paper["pass_score"] or 60
            online = False
        # 得分判定统一为：得分率 = 得分/满分，再与通过线(pass_score)比较
        sr, passed = self._score_rate_pass(round(total, 2), full_score, ps, ratio, online=online)
        # 回写 score_rate / full_score，否则简答题判分后「得分率%」仍是判分前的旧值，
        # 导致个人视图与管理端显示/判定不一致。
        self.conn.execute(
            "UPDATE exam_attempts SET manual_score=?, total_score=?, score_rate=?, full_score=?, passed=?, status='graded' WHERE attempt_id=?",
            (round(manual, 2), round(total, 2), round(sr, 4) if sr is not None else None, full_score, passed, attempt_id))
        self.conn.commit()
        self._sync_attempt_to_results(attempt_id)
        return self.get_attempt(attempt_id)

    # ---- 种子数据 ----
    def seed_if_empty(self) -> None:
        rep_count = self._row(self.conn.execute("SELECT COUNT(*) c FROM cs_reps"))["c"]
        if rep_count > 0:
            return
        reps = [
            ("CS001", "Irish", "2026-01-05"),
            ("CS002", "Jeniffer", "2026-01-08"),
            ("CS003", "Kate", "2026-02-01"),
            ("CS004", "Fema", "2026-02-15"),
            ("CS005", "Meliza", "2026-03-01"),
        ]
        for rid, name, hd in reps:
            self.create_rep({"rep_id": rid, "name": name, "hire_date": hd})

        sessions = [
            ("入职培训考试", "2026-03-Batch1", "2026-03-10", 60, "onboarding"),
            ("产品知识月考", "2026-04-Batch1", "2026-04-12", 60, "monthly"),
            ("服务流程季度考", "2026-05-Batch1", "2026-05-15", 70, "quarterly"),
        ]
        import random
        rng = random.Random(42)
        for name, batch, date, pscore, etype in sessions:
            results = []
            for rid, name_, hd in reps:
                subj = {
                    "Q1": rng.randint(55, 98),
                    "Q2": rng.randint(60, 99),
                    "Q3": rng.randint(50, 95),
                }
                results.append({
                    "rep_id": rid, "name": name_,
                    "subjects": subj,
                })
            self.create_session(
                {"exam_name": name, "batch": batch, "exam_date": date,
                 "pass_score": pscore, "exam_type": etype, "note": "种子数据演示"},
                results=results)
        # 给部分客服增加第二次考试以展示个人趋势
        extra = [
            ("CS001", "2026-05-Batch2", "2026-05-20", {"Q1": 92, "Q2": 95, "Q3": 88}),
            ("CS002", "2026-05-Batch2", "2026-05-20", {"Q1": 78, "Q2": 82, "Q3": 75}),
        ]
        sid = self.conn.execute(
            "SELECT session_id FROM exam_sessions WHERE batch='2026-05-Batch2'").fetchone()
        if not sid:
            sess = self.create_session(
                {"exam_name": "服务流程季度考(补考)", "batch": "2026-05-Batch2",
                 "exam_date": "2026-05-20", "pass_score": 70, "exam_type": "quarterly"})
            sid = sess["session_id"]
        for rid, batch, date, subj in extra:
            rep = self._row(self.conn.execute("SELECT name FROM cs_reps WHERE rep_id=?", (rid,)))
            self.create_result({"session_id": sid, "rep_id": rid,
                                 "name": rep["name"], "subjects": subj})




# ----------------------------------------------------------------------------
# 第二阶段：云端存储 占位实现（演示可切换）
# ----------------------------------------------------------------------------
class CloudStorage(BaseStorage):
    """
    第二阶段（云端）实现占位。

    迁移策略（二选一或结合）：
      A. 自建 REST 服务：把 SQLiteStorage 的 SQL 换成对远端 API 的 HTTP 调用，
         每个方法对应一个端点，例如 GET /reps、POST /results。
      B. 文件型云盘（如 Google Drive）：把整库序列化为 JSON / 把 exam.db 上传，
         读时下载、写后回传（适合低频写入场景）。

    本类仅演示接口骨架，未实现具体网络调用。切换时：
        设置环境变量 STORAGE_BACKEND=cloud 即可由 get_storage() 返回本类。
    由于所有业务方法签名与 SQLiteStorage 完全一致，server.py 与前端无需改动。
    """

    def __init__(self, endpoint: str = None):
        self.endpoint = endpoint or os.environ.get("CLOUD_ENDPOINT", "")
        if not self.endpoint:
            raise RuntimeError(
                "CloudStorage 需要设置 CLOUD_ENDPOINT（远端 REST 地址或云盘根路径）")

    def _todo(self, name: str):
        raise NotImplementedError(
            f"[CloudStorage] 方法 {name} 尚未实现 —— 请在第二阶段接入云端后端 "
            f"({self.endpoint})。接口签名已与 SQLiteStorage 对齐，实现后直接替换即可。")

    def init(self): self._todo("init")
    def seed_if_empty(self): self._todo("seed_if_empty")
    def list_reps(self): self._todo("list_reps")
    def get_rep_by_name(self, name): self._todo("get_rep_by_name")
    def create_rep(self, data): self._todo("create_rep")
    def update_rep(self, rep_id, data): self._todo("update_rep")
    def delete_rep(self, rep_id): self._todo("delete_rep")
    def list_sessions(self): self._todo("list_sessions")
    def get_session(self, session_id): self._todo("get_session")
    def create_session(self, data, results=None): self._todo("create_session")
    def update_session(self, session_id, data): self._todo("update_session")
    def delete_session(self, session_id): self._todo("delete_session")
    def bulk_delete_sessions(self, ids): self._todo("bulk_delete_sessions")
    def list_results(self, filters=None): self._todo("list_results")
    def create_result(self, data): self._todo("create_result")
    def update_result(self, result_id, data): self._todo("update_result")
    def delete_result(self, result_id): self._todo("delete_result")
    def bulk_delete_results(self, ids): self._todo("bulk_delete_results")
    def reset_all_data(self, scope="all"): self._todo("reset_all_data")
    def import_results_excel(self, path, exam_name=None, batch=None, exam_date=None,
                             pass_ratio=0.7, full_score=None, pass_score=None,
                             orig_filename=None):
        self._todo("import_results_excel")
    def import_questions_ppt(self, path, exam_name, dim_ids=None, orig_filename=None):
        self._todo("import_questions_ppt")
    def list_question_banks(self): self._todo("list_question_banks")
    def delete_question_bank(self, exam_name): self._todo("delete_question_bank")
    def bulk_delete_question_banks(self, names): self._todo("bulk_delete_question_banks")
    def get_question_by_exam_q(self, exam_name, q_number):
        self._todo("get_question_by_exam_q")
    def list_dimensions(self): self._todo("list_dimensions")
    def create_dimension(self, data): self._todo("create_dimension")
    def update_dimension(self, dim_id, data): self._todo("update_dimension")
    def delete_dimension(self, dim_id): self._todo("delete_dimension")
    def bulk_delete_dimensions(self, ids): self._todo("bulk_delete_dimensions")
    def set_question_dimension(self, exam_name, q_number, dim_id, max_score=None): self._todo("set_question_dimension")
    def delete_question_dimension(self, exam_name, q_number, dim_id): self._todo("delete_question_dimension")
    def get_exam_question_dimensions(self, exam_name): self._todo("get_exam_question_dimensions")
    def exam_dimension_distribution(self, exam_name): self._todo("exam_dimension_distribution")
    def rep_dimension_weakness(self, rep_id, session_ids=None): self._todo("rep_dimension_weakness")
    def individual_view(self, rep_id): self._todo("individual_view")
    def batch_view(self, session_id): self._todo("batch_view")
    def period_view(self, start, end): self._todo("period_view")
    def overview(self): self._todo("overview")

    # ---- 积分系统（占位桩，接口与 SQLiteStorage 对齐）----
    def get_points_rules(self): self._todo("get_points_rules")
    def get_points_threshold(self): self._todo("get_points_threshold")
    def pass_points_for(self, score_rate, passed=False): self._todo("pass_points_for")
    def award_points(self, rep_id, rule_key, ref_type, ref_id, note="", delta=None): self._todo("award_points")
    def points_summary(self): self._todo("points_summary")
    def points_year_summary(self, year=None): self._todo("points_year_summary")
    def list_point_years(self): self._todo("list_point_years")
    def rep_points(self, rep_id, year=None): self._todo("rep_points")
    def rep_points_log(self, rep_id, year=None): self._todo("rep_points_log")
    def update_points_config(self, rules=None, threshold=None): self._todo("update_points_config")


# ----------------------------------------------------------------------------
# 工厂：通过配置切换后端，业务代码零改动
# ----------------------------------------------------------------------------
# ----------------------------------------------------------------------------
# Neon / Postgres 实现（云端生产后端）
# ----------------------------------------------------------------------------
# 通过 _PgConn 代理，把 SQLiteStorage 的 SQL 实时翻译成 Postgres，
# 从而继承其全部 254 个业务方法，无需重写。
PG_SCHEMA = """
CREATE TABLE IF NOT EXISTS cs_reps (
    rep_id     TEXT PRIMARY KEY,
    name       TEXT NOT NULL,
    hire_date  TEXT,
    leave_date TEXT,
    status     TEXT DEFAULT 'active',
    stage      TEXT DEFAULT '新人',
    login_id   TEXT,
    position   TEXT,
    channel    TEXT,
    created_at TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS'))
);

CREATE TABLE IF NOT EXISTS knowledge_dimensions (
    dim_id      SERIAL PRIMARY KEY,
    name_cn     TEXT NOT NULL,
    name_en     TEXT,
    description TEXT,
    created_at  TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS'))
);

CREATE TABLE IF NOT EXISTS questions (
    question_id SERIAL PRIMARY KEY,
    q_type      TEXT NOT NULL DEFAULT 'single',
    category    TEXT,
    content     TEXT NOT NULL,
    options     TEXT,
    answer      TEXT,
    dim_id      INTEGER,
    score       REAL DEFAULT 5,
    explanation TEXT,
    source_exam TEXT,
    created_at  TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS')),
    FOREIGN KEY (dim_id) REFERENCES knowledge_dimensions(dim_id)
);

CREATE TABLE IF NOT EXISTS exam_sessions (
    session_id  SERIAL PRIMARY KEY,
    exam_name   TEXT NOT NULL,
    batch       TEXT NOT NULL,
    exam_date   TEXT NOT NULL,
    pass_score  REAL DEFAULT 60,
    exam_type   TEXT DEFAULT 'onboarding',
    cycle_tag   TEXT,
    note        TEXT,
    created_at  TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS'))
);

CREATE TABLE IF NOT EXISTS exam_papers (
    paper_id    SERIAL PRIMARY KEY,
    title       TEXT NOT NULL,
    batch       TEXT DEFAULT 'online',
    exam_type   TEXT DEFAULT 'onboarding',
    status      TEXT DEFAULT 'draft',
    duration_min INTEGER DEFAULT 0,
    open_at     TEXT,
    close_at    TEXT,
    pass_score  REAL DEFAULT 60,
    created_by  TEXT,
    created_at  TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS'))
);

CREATE TABLE IF NOT EXISTS exam_results (
    result_id  SERIAL PRIMARY KEY,
    session_id INTEGER NOT NULL,
    rep_id     TEXT NOT NULL,
    name       TEXT NOT NULL,
    subjects   TEXT,
    total      REAL,
    passed     INTEGER,
    score_rate REAL,
    full_score REAL,
    created_at TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS')),
    FOREIGN KEY (session_id) REFERENCES exam_sessions(session_id),
    FOREIGN KEY (rep_id) REFERENCES cs_reps(rep_id)
);

CREATE TABLE IF NOT EXISTS exam_question_dimensions (
    exam_name   TEXT NOT NULL,
    q_number    TEXT NOT NULL,
    dim_id      INTEGER NOT NULL,
    max_score   REAL,
    PRIMARY KEY (exam_name, q_number, dim_id),
    FOREIGN KEY (dim_id) REFERENCES knowledge_dimensions(dim_id)
);

CREATE TABLE IF NOT EXISTS exam_question_bank (
    exam_name   TEXT NOT NULL,
    q_number    TEXT NOT NULL,
    question_id INTEGER NOT NULL,
    seq         INTEGER,
    PRIMARY KEY (exam_name, q_number),
    FOREIGN KEY (question_id) REFERENCES questions(question_id)
);

CREATE TABLE IF NOT EXISTS question_bank_meta (
    exam_name     TEXT PRIMARY KEY,
    orig_filename TEXT,
    uploaded_at   TEXT,
    q_count       INTEGER,
    dim_ids       TEXT
);

CREATE TABLE IF NOT EXISTS accounts (
    rep_id        TEXT PRIMARY KEY,
    login_id      TEXT UNIQUE,
    password_hash TEXT,
    role          TEXT DEFAULT 'csr',
    created_at    TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS')),
    FOREIGN KEY (rep_id) REFERENCES cs_reps(rep_id)
);

CREATE TABLE IF NOT EXISTS question_attachments (
    att_id      SERIAL PRIMARY KEY,
    question_id INTEGER NOT NULL,
    seq         INTEGER DEFAULT 0,
    filename    TEXT,
    mime        TEXT,
    stored_path TEXT,
    data        BYTEA,
    created_at  TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS')),
    FOREIGN KEY (question_id) REFERENCES questions(question_id)
);

CREATE TABLE IF NOT EXISTS exam_assignments (
    paper_id    INTEGER NOT NULL,
    rep_id      TEXT NOT NULL,
    open_at     TEXT,
    due_at      TEXT,
    created_at  TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS')),
    PRIMARY KEY (paper_id, rep_id),
    FOREIGN KEY (paper_id) REFERENCES exam_papers(paper_id),
    FOREIGN KEY (rep_id) REFERENCES cs_reps(rep_id)
);

CREATE TABLE IF NOT EXISTS paper_questions (
    paper_id    INTEGER NOT NULL,
    question_id INTEGER NOT NULL,
    seq         INTEGER,
    score       REAL,
    PRIMARY KEY (paper_id, question_id),
    FOREIGN KEY (paper_id) REFERENCES exam_papers(paper_id),
    FOREIGN KEY (question_id) REFERENCES questions(question_id)
);

CREATE TABLE IF NOT EXISTS exam_attempts (
    attempt_id  SERIAL PRIMARY KEY,
    paper_id    INTEGER NOT NULL,
    rep_id      TEXT NOT NULL,
    start_time  TEXT,
    submit_time TEXT,
    status      TEXT DEFAULT 'in_progress',
    auto_score  REAL DEFAULT 0,
    manual_score REAL DEFAULT 0,
    total_score REAL DEFAULT 0,
    passed      INTEGER DEFAULT 0,
    score_rate  REAL,
    full_score  REAL,
    created_at  TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS')),
    FOREIGN KEY (paper_id) REFERENCES exam_papers(paper_id),
    FOREIGN KEY (rep_id) REFERENCES cs_reps(rep_id)
);

CREATE TABLE IF NOT EXISTS exam_answers (
    attempt_id  INTEGER NOT NULL,
    question_id INTEGER NOT NULL,
    answer      TEXT,
    is_correct  INTEGER,
    score       REAL DEFAULT 0,
    graded_by   TEXT,
    graded_at   TEXT,
    PRIMARY KEY (attempt_id, question_id),
    FOREIGN KEY (attempt_id) REFERENCES exam_attempts(attempt_id),
    FOREIGN KEY (question_id) REFERENCES questions(question_id)
);

CREATE TABLE IF NOT EXISTS system_config (
    key        TEXT PRIMARY KEY,
    value      TEXT,
    updated_at TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS'))
);

CREATE TABLE IF NOT EXISTS points_account (
    rep_id TEXT PRIMARY KEY, total INTEGER DEFAULT 0,
    updated_at TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS'))
);

CREATE TABLE IF NOT EXISTS points_log (
    log_id SERIAL PRIMARY KEY, rep_id TEXT, rule_key TEXT,
    delta INTEGER, ref_type TEXT, ref_id TEXT, note TEXT,
    created_at TEXT, year INTEGER, quarter INTEGER
);

CREATE TABLE IF NOT EXISTS learning_materials (
    material_id SERIAL PRIMARY KEY, title TEXT, mtype TEXT,
    dim_id INTEGER, content TEXT, file_path TEXT, url TEXT, link_kind TEXT,
    created_at TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS'))
);

CREATE TABLE IF NOT EXISTS material_dimensions (
    material_id INTEGER NOT NULL, dim_id INTEGER NOT NULL,
    PRIMARY KEY (material_id, dim_id)
);

CREATE TABLE IF NOT EXISTS study_records (
    record_id SERIAL PRIMARY KEY, rep_id TEXT, material_id INTEGER,
    status TEXT DEFAULT 'opened', progress INTEGER DEFAULT 0,
    created_at TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS')),
    UNIQUE (rep_id, material_id)
);

CREATE TABLE IF NOT EXISTS mini_quiz (
    quiz_id SERIAL PRIMARY KEY, rep_id TEXT, dim_id INTEGER,
    question_ids TEXT, passed INTEGER DEFAULT 0,
    created_at TEXT DEFAULT (to_char(now(), 'YYYY-MM-DD HH24:MI:SS'))
);

CREATE INDEX IF NOT EXISTS idx_results_session ON exam_results(session_id);
CREATE INDEX IF NOT EXISTS idx_results_rep     ON exam_results(rep_id);
CREATE INDEX IF NOT EXISTS idx_paper_q_paper  ON paper_questions(paper_id);
CREATE INDEX IF NOT EXISTS idx_attempt_paper  ON exam_attempts(paper_id);
CREATE INDEX IF NOT EXISTS idx_attempt_rep    ON exam_attempts(rep_id);
"""

# INSERT OR REPLACE 的冲突目标（主键列），用于翻译成 ON CONFLICT DO UPDATE
_PG_REPLACE_PK = {
    "exam_assignments": ("paper_id", "rep_id"),
    "accounts": ("rep_id",),
    "cs_reps": ("rep_id",),
    "exam_question_dimensions": ("exam_name", "q_number", "dim_id"),
    "study_records": ("rep_id", "material_id"),
    "exam_question_bank": ("exam_name", "q_number"),
    "question_bank_meta": ("exam_name",),
}
# 单 SERIAL 主键表：INSERT 后追加 RETURNING <pk> 以精确模拟 lastrowid
_PG_SERIAL_PK = {
    "exam_sessions": "session_id",
    "questions": "question_id",
    "question_attachments": "att_id",
    "exam_papers": "paper_id",
    "exam_attempts": "attempt_id",
    "exam_results": "result_id",
    "points_log": "log_id",
    "learning_materials": "material_id",
    "mini_quiz": "quiz_id",
    "knowledge_dimensions": "dim_id",
}


def _pg_translate(sql: str) -> str:
    """把 SQLite 风格 SQL 翻成 Postgres 风格（仅做安全、确定的改写）。"""
    s = sql.replace("?", "%s")
    # INSERT OR IGNORE -> ... VALUES (...) ON CONFLICT DO NOTHING
    s = re.sub(
        r"INSERT OR IGNORE INTO (\w+)\s*\(([^)]*)\)\s*VALUES\s*\(([^)]*)\)",
        r"INSERT INTO \1 (\2) VALUES (\3) ON CONFLICT DO NOTHING",
        s, flags=re.IGNORECASE)
    # INSERT OR REPLACE -> ... VALUES (...) ON CONFLICT (pk) DO UPDATE SET ...
    m = re.search(
        r"INSERT OR REPLACE INTO (\w+)\s*\(([^)]*)\)\s*VALUES\s*\(([^)]*)\)",
        s, flags=re.IGNORECASE)
    if m:
        tbl, colstr, valstr = m.groups()
        pk = _PG_REPLACE_PK.get(tbl)
        if pk:
            cols = [c.strip() for c in colstr.split(",")]
            nonpk = [c for c in cols if c not in pk]
            setcl = ", ".join(f"{c}=EXCLUDED.{c}" for c in nonpk)
            s = (f"INSERT INTO {tbl} ({colstr}) VALUES ({valstr}) "
                 f"ON CONFLICT ({', '.join(pk)}) DO UPDATE SET {setcl}")
    # SQLite 日期函数 strftime('fmt', col) -> TO_CHAR(col::timestamp, 'PGFMT')
    def _strftime_repl(mm):
        fmt = mm.group(1)
        col = mm.group(2)
        pg = (fmt.replace("%Y", "YYYY").replace("%m", "MM").replace("%d", "DD")
                  .replace("%H", "HH24").replace("%M", "MI").replace("%S", "SS"))
        return f"TO_CHAR({col}::timestamp, '{pg}')"
    s = re.sub(r"strftime\('([^']+)',\s*([^)]+)\)", _strftime_repl, s)
    return s


class _PgCursor:
    """包装 psycopg2 cursor：行按 dict 访问、INSERT 后精确模拟 lastrowid。"""

    def __init__(self, cur, conn):
        self._cur = cur
        self._conn = conn
        self._lastval = None

    def execute(self, sql, params=None):
        sql2 = _pg_translate(sql)
        pk = None
        mm = re.match(r"\s*INSERT INTO (\w+)", sql2, flags=re.IGNORECASE)
        if mm:
            pk = _PG_SERIAL_PK.get(mm.group(1))
        if pk:
            sql2 = sql2.rstrip().rstrip(";") + f" RETURNING {pk}"
        # SQLite 把 Python bool 存成 0/1；Postgres 的 INTEGER 列不接受布尔字面量，
        # 故在发送前把 bool 统一转成 int，sqlite 端不受影响（int 仍是 int）。
        if params is not None:
            if isinstance(params, (list, tuple)):
                params = tuple(int(p) if isinstance(p, bool) else p for p in params)
            elif isinstance(params, dict):
                params = {k: (int(v) if isinstance(v, bool) else v) for k, v in params.items()}
        self._cur.execute(sql2, params)
        if pk:
            row = self._cur.fetchone()
            self._lastval = row[0] if row else None
        return self

    def executemany(self, sql, params_seq):
        self._cur.executemany(_pg_translate(sql), params_seq)
        return self

    def fetchone(self):
        r = self._cur.fetchone()
        return dict(r) if r is not None else None

    def fetchall(self):
        return [dict(r) for r in self._cur.fetchall()]

    @property
    def lastrowid(self):
        return self._lastval

    @property
    def rowcount(self):
        return self._cur.rowcount

    def close(self):
        try:
            self._cur.close()
        except Exception:
            pass


class _PgConn:
    """代理连接：把所有执行请求翻译后交给 psycopg2。"""

    def __init__(self, pg):
        self._pg = pg

    def cursor(self):
        return _PgCursor(self._pg.cursor(), self._pg)

    def execute(self, sql, params=None):
        return _PgCursor(self._pg.cursor(), self._pg).execute(sql, params)

    def commit(self):
        self._pg.commit()

    def rollback(self):
        self._pg.rollback()

    def close(self):
        try:
            self._pg.close()
        except Exception:
            pass


class PostgresStorage(SQLiteStorage):
    """Neon/Postgres 实现：继承 SQLiteStorage 全部业务方法，仅替换连接层与 DDL。"""

    def __init__(self, database_url: str = None):
        self.database_url = database_url or os.environ.get("DATABASE_URL", "")
        if not self.database_url:
            raise RuntimeError("PostgresStorage 需要设置 DATABASE_URL 环境变量（Neon 连接串）")
        if not _HAS_PG:
            raise RuntimeError("未安装 psycopg2-binary，无法使用 Postgres 后端")
        self._local = threading.local()

    def _new_conn(self):
        pg = psycopg2.connect(self.database_url, cursor_factory=DictCursor)
        pg.autocommit = False
        return pg

    @property
    def conn(self):
        try:
            from flask import g
            if not hasattr(g, "pg_conn") or g.pg_conn is None:
                g.pg_conn = self._new_conn()
            return _PgConn(g.pg_conn)
        except RuntimeError:
            t = self._local
            if not hasattr(t, "conn") or t.conn is None:
                t.conn = self._new_conn()
            return _PgConn(t.conn)

    def init(self) -> None:
        pg = self._new_conn()
        cur = pg.cursor()
        for stmt in PG_SCHEMA.split(";"):
            if stmt.strip():
                cur.execute(stmt)
        # 迁移：已存在的库补 data 列（图片改存数据库，跨部署不丢）
        cur.execute("ALTER TABLE question_attachments ADD COLUMN IF NOT EXISTS data BYTEA")
        # 清理迁移前遗留的、只有磁盘路径而无二进制数据的孤儿附件（图片已无法读取）
        cur.execute("DELETE FROM question_attachments WHERE data IS NULL")
        pg.commit()
        cur.close()
        pg.close()

    def seed_if_empty(self) -> None:
        # 生产环境：不灌 demo 客服/成绩（考题+资料库通过迁移脚本导入）。
        # 仅确保系统设置与管理员密码哈希存在，使管理员可登录。
        self._ensure_config()
        self._seed_default_dimensions_pg()

    def _ensure_config(self) -> None:
        def set_cfg(key, value):
            self.conn.execute(
                "INSERT INTO system_config (key, value, updated_at) VALUES (?,?,?) "
                "ON CONFLICT (key) DO UPDATE SET value=EXCLUDED.value, updated_at=EXCLUDED.updated_at",
                (key, value, _now()))

        set_cfg("pass_line_ratio", os.environ.get("PASS_LINE_RATIO", "0.88"))
        set_cfg("points_rules", os.environ.get("POINTS_RULES", json.dumps(
            {"participate": 10, "pass": {"0.8": 15, "0.9": 20, "0.95": 30},
             "mini_quiz": 10, "material": 5})))
        set_cfg("points_period", os.environ.get("POINTS_PERIOD", "quarter"))
        set_cfg("points_period_target", os.environ.get("POINTS_PERIOD_TARGET", "0"))
        admin_pw = os.environ.get("ADMIN_PW", "admin123")
        # 仅在首次（哈希不存在）时写入；不覆盖管理员在界面修改过的密码，
        # 否则每次冷启动（如 Render 免费版休眠后唤醒）都会把密码重置回 ADMIN_PW。
        if not self._row(self.conn.execute("SELECT 1 FROM system_config WHERE key='admin_password_hash'")):
            set_cfg("admin_password_hash", _hash_pw(admin_pw))
        self.conn.commit()

    def _seed_default_dimensions_pg(self) -> None:
        c = self._row(self.conn.execute("SELECT COUNT(*) c FROM knowledge_dimensions"))["c"]
        if c == 0:
            for cn, en in DEFAULT_DIMS:
                self.conn.execute(
                    "INSERT INTO knowledge_dimensions (name_cn, name_en, created_at) VALUES (?,?,?)",
                    (cn, en, _now()))
            self.conn.commit()


def get_storage() -> BaseStorage:
    if os.environ.get("DATABASE_URL"):
        if not _HAS_PG:
            raise RuntimeError("未安装 psycopg2，无法使用 Postgres 后端")
        return PostgresStorage()
    if BACKEND == "cloud":
        return CloudStorage()
    return SQLiteStorage(DB_PATH)


if __name__ == "__main__":
    s = get_storage()
    s.init()
    s.seed_if_empty()
    print("存储后端:", BACKEND)
    print("总览:", s.overview())
